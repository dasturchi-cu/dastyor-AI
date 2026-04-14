"""WebApp JSON APIs: auth, profile, messaging, transliteration, translate, spellcheck, objective."""
from __future__ import annotations

import asyncio
import html as html_lib
import io
import hashlib
import logging
import os
import tempfile
from typing import Optional

from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import StreamingResponse
from telegram import InputFile

from backend.dependencies import get_ptb_application
from backend.schemas.webapp import (
    AuthRequest,
    NotifyRequest,
    ObjectiveRequest,
    SpellcheckRequest,
    SupportRequest,
    TranslateAutoRequest,
    TranslateRequest,
    TranslitRequest,
)
from backend.services.redis_json_cache import (
    api_me_cache_key,
    redis_cache_delete,
    redis_cache_get_json,
    redis_cache_set_json,
)
from backend.services.spellcheck_cache import spellcheck_cache_get, spellcheck_cache_key, spellcheck_cache_set
from backend.services.upload_io import EmptyUploadError, UploadTooLargeError, read_upload_limited
from backend.services.user_resolve import resolve_telegram_uid
from backend.services.web_quota import web_quota_commit_success, web_quota_consume_or_raise
from backend.settings import get_settings
from backend.web_constants import (
    BOT_USERNAME,
    NOTIFY_MAX_CHARS,
    SPELLCHECK_MAX_CHARS,
    TRANSLATE_MAX_CHARS,
    TRANSLIT_MAX_CHARS,
    WEBAPP_BASE,
    WEBAPP_VERSION,
)
from backend.services.auto_script import auto_cyrillic_latin
from backend.services.supabase_storage import create_signed_url, upload_bytes_for_user
from backend.services.cv_parser import parse_cv_text

logger = logging.getLogger(__name__)
SPELLCHECK_FILE_MAX_CHARS = int(os.getenv("SPELLCHECK_FILE_MAX_CHARS", "150000"))
TRANSLATE_AUTO_CACHE_TTL_SECONDS = int(os.getenv("TRANSLATE_AUTO_CACHE_TTL_SECONDS", "86400"))

router = APIRouter(tags=["web-api"])


def _resolve_web_uid_optional(telegram_id: Optional[int], token: Optional[str]) -> Optional[int]:
    uid = resolve_telegram_uid(str(telegram_id) if telegram_id is not None else None, token)
    return int(uid) if uid else None


def _build_pdf_from_text(text: str, out_path: str) -> None:
    """Create a simple UTF-8 safe-ish PDF from plain text."""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    # Try to register a Unicode font for Uzbek/Cyrillic support.
    try:
        font_path = os.path.join(os.getcwd(), "assets", "fonts", "DejaVuSans.ttf")
        if os.path.exists(font_path):
            pdfmetrics.registerFont(TTFont("DejaVuSans", font_path))
            font_name = "DejaVuSans"
        else:
            font_name = "Helvetica"
    except Exception:
        font_name = "Helvetica"

    c = canvas.Canvas(out_path, pagesize=A4)
    width, height = A4
    left = 40
    top = height - 50
    line_h = 15
    y = top
    c.setFont(font_name, 11)
    for raw_line in (text or "").splitlines() or [""]:
        line = (raw_line or "").replace("\t", "    ")
        if y < 50:
            c.showPage()
            c.setFont(font_name, 11)
            y = top
        c.drawString(left, y, line[:1800])
        y -= line_h
    c.save()


@router.post("/api/auth")
async def api_auth(req: AuthRequest):
    try:
        from bot.services.session_service import create_session
        from bot.services.user_service import track_user_activity

        token = create_session(
            telegram_id=req.telegram_id,
            first_name=req.first_name,
            username=req.username,
            photo_url=req.photo_url,
        )

        class _FakeUser:
            id = req.telegram_id
            first_name = req.first_name
            username = req.username

        track_user_activity(
            _FakeUser(), command="web_auth", chat_id=int(req.telegram_id)
        )
        await redis_cache_delete(api_me_cache_key(str(req.telegram_id)))
        return {"ok": True, "token": token, "telegram_id": req.telegram_id}
    except Exception as e:
        logger.error("/api/auth error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=str(e)[:200])


def _build_api_me_payload(uid: str) -> dict:
    """Sync DB/file work — run in a thread pool so the event loop stays responsive."""
    from bot.services.session_service import get_session_by_telegram_id
    from bot.services.settings_service import is_premium
    from bot.services.user_service import get_user_profile
    from bot.services.usage_tracker import get_tariff_snapshot

    profile = get_user_profile(uid) or {}
    session = get_session_by_telegram_id(uid) or {}
    tariff = get_tariff_snapshot(int(uid))

    return {
        "ok": True,
        "telegram_id": uid,
        "first_name": session.get("first_name", profile.get("first_name", "")),
        "username": session.get("username", profile.get("username", "")),
        "photo_url": session.get("photo_url", ""),
        "is_premium": is_premium(int(uid)),
        "files_processed": profile.get("files_processed", 0),
        "joined_at": profile.get("joined_at", ""),
        "last_active": profile.get("last_active", ""),
        # Referral status (used by premium.html banner and pricing)
        "referred_by": profile.get("referred_by"),
        "referrals_count": int(profile.get("referrals_count") or 0),
        "referral_discount_percent": int(profile.get("referral_discount_percent") or 0),
        "referral_discount_active": bool(profile.get("referral_discount_active")),
        "referral_discount_expires_at": profile.get("referral_discount_expires_at"),
        # Per-plan referral discount (Standard/Premium separately)
        "referral_discount_standard_active": bool(profile.get("referral_discount_standard_active")),
        "referral_discount_premium_active": bool(profile.get("referral_discount_premium_active")),
        "referral_discount_standard_consumed_at": profile.get("referral_discount_standard_consumed_at"),
        "referral_discount_premium_consumed_at": profile.get("referral_discount_premium_consumed_at"),
        **tariff,
    }


@router.get("/api/me")
async def api_me(
    token: Optional[str] = Query(None),
    telegram_id: Optional[str] = Query(None),
):
    uid = resolve_telegram_uid(telegram_id, token)
    if not uid:
        raise HTTPException(status_code=401, detail="Foydalanuvchi aniqlanmadi")

    cache_ttl = int(os.getenv("API_ME_CACHE_TTL_SECONDS", "15") or "15")
    ck = api_me_cache_key(uid)
    if cache_ttl > 0:
        hit = await redis_cache_get_json(ck)
        if hit is not None:
            return hit

    body = await asyncio.to_thread(_build_api_me_payload, uid)
    if cache_ttl > 0:
        await redis_cache_set_json(ck, body, cache_ttl)
    return body


@router.post("/api/translit")
async def api_translit(req: TranslitRequest):
    if not req.text or not req.text.strip():
        raise HTTPException(status_code=400, detail="Matn bo'sh bo'lishi mumkin emas")
    if len(req.text) > TRANSLIT_MAX_CHARS:
        raise HTTPException(status_code=400, detail="Matn 50 000 belgidan oshmasligi kerak")

    valid = {"krill_to_lotin", "lotin_to_krill"}
    if req.direction not in valid:
        raise HTTPException(status_code=400, detail=f"Noto'g'ri yo'nalish: {req.direction}")

    uid = _resolve_web_uid_optional(req.telegram_id, req.token)

    quota = None
    consumed = False
    try:
        from bot.services.transliterate_service import transliterate

        if uid:
            from bot.services.plan_limits import CAT_TRANSLIT

            consumed = bool(web_quota_consume_or_raise(int(uid), CAT_TRANSLIT))

        result = transliterate(req.text, req.direction)  # type: ignore[arg-type]
        if uid:
            from bot.services.plan_limits import CAT_TRANSLIT

            from bot.services.plan_limits import category_quota_for_response
            from bot.services.user_service import record_service_completion

            record_service_completion(int(uid), CAT_TRANSLIT, "Web translit", skip_quota=bool(consumed))
            quota = category_quota_for_response(int(uid), CAT_TRANSLIT)
            try:
                from bot.utils.action_logger import log_action_fire_and_forget

                log_action_fire_and_forget(
                    telegram_id=int(uid),
                    username=None,
                    action_type="TRANSLIT",
                    details=f"web:{req.direction}",
                    metadata={"direction": req.direction, "chars": len(req.text or "")},
                )
            except Exception:
                pass
        return {"ok": True, "result": result, "direction": req.direction, "quota": quota}
    except HTTPException:
        raise
    except Exception as e:
        logger.error("/api/translit error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=str(e)[:200])


@router.post("/api/translit_auto")
async def api_translit_auto(req: SpellcheckRequest):
    """
    Auto Cyrillic ↔ Latin:
      - Latin dominates -> convert to Cyrillic
      - Cyrillic dominates -> convert to Latin
    No user input required.
    """
    if not req.text or not req.text.strip():
        raise HTTPException(status_code=400, detail="Matn bo'sh bo'lishi mumkin emas")
    if len(req.text) > TRANSLIT_MAX_CHARS:
        raise HTTPException(status_code=400, detail="Matn 50 000 belgidan oshmasligi kerak")

    uid = _resolve_web_uid_optional(req.telegram_id, req.token)
    try:
        res = auto_cyrillic_latin(req.text)
        quota = None
        consumed = False
        if uid and res.direction != "none":
            from bot.services.plan_limits import CAT_TRANSLIT

            consumed = bool(web_quota_consume_or_raise(int(uid), CAT_TRANSLIT))
            from bot.services.plan_limits import category_quota_for_response
            from bot.services.user_service import record_service_completion

            record_service_completion(int(uid), CAT_TRANSLIT, "Web translit auto", skip_quota=bool(consumed))
            quota = category_quota_for_response(int(uid), CAT_TRANSLIT)
            try:
                from bot.utils.action_logger import log_action_fire_and_forget

                log_action_fire_and_forget(
                    telegram_id=int(uid),
                    username=None,
                    action_type="TRANSLIT_AUTO",
                    details="web:auto",
                    metadata={
                        "detected": res.detected,
                        "direction": res.direction,
                        "chars": len(req.text or ""),
                        "latin": (res.stats or {}).get("latin"),
                        "cyrillic": (res.stats or {}).get("cyrillic"),
                    },
                )
            except Exception:
                pass
        return {
            "ok": True,
            "result": res.result,
            "direction": res.direction,
            "detected": res.detected,
            "stats": res.stats,
            "quota": quota,
        }
    except Exception as e:
        logger.error("/api/translit_auto error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=str(e)[:200])


@router.get("/api/bot-link")
async def api_bot_link(
    action: str = Query(..., description="cv | obyektivka | ocr | pdf | translit | translate"),
    telegram_id: Optional[str] = Query(None),
):
    page_map = {
        "cv": "cv.html",
        "obyektivka": "obyektivka.html",
        "ocr": "ocr.html",
        "pdf": "img2pdf.html",
        "translit": "translit.html",
        "translate": "translate.html",
        "premium": "premium.html",
    }
    page = page_map.get(action)
    if not page:
        raise HTTPException(status_code=400, detail=f"Noma'lum action: {action}")

    bot_link = f"https://t.me/{BOT_USERNAME}?start={action}"
    tid_param = f"?telegram_id={telegram_id}" if telegram_id else ""
    v_param = f"{'&' if tid_param else '?'}v={WEBAPP_VERSION}"
    webapp_url = f"{WEBAPP_BASE}/{page}{tid_param}{v_param}"
    return {"ok": True, "bot_link": bot_link, "webapp_url": webapp_url, "action": action}


@router.get("/api/stats")
async def api_stats(
    token: Optional[str] = Query(None),
    telegram_id: Optional[str] = Query(None),
):
    uid = resolve_telegram_uid(telegram_id, token)
    if not uid:
        raise HTTPException(status_code=401, detail="Foydalanuvchi aniqlanmadi")

    from bot.services.settings_service import is_premium
    from bot.services.user_service import get_user_profile

    profile = get_user_profile(uid) or {}
    premium = is_premium(int(uid))

    return {
        "ok": True,
        "telegram_id": uid,
        "is_premium": premium,
        "files_processed": profile.get("files_processed", 0),
        "sessions": profile.get("sessions", 1),
        "last_service": profile.get("last_service", ""),
        "last_active": profile.get("last_active", ""),
        "joined_at": profile.get("joined_at", ""),
    }


@router.get("/api/referrals")
async def api_referrals(
    token: Optional[str] = Query(None),
    telegram_id: Optional[str] = Query(None),
    limit: int = Query(20, ge=1, le=50),
):
    """
    Return referral progress and invitees list for the current user.
    """
    uid = resolve_telegram_uid(telegram_id, token)
    if not uid:
        raise HTTPException(status_code=401, detail="Foydalanuvchi aniqlanmadi")

    # Prefer Supabase; fallback: no list
    from bot.services.user_service import get_user_profile

    profile = get_user_profile(uid) or {}
    cnt = int(profile.get("referrals_count") or 0)
    active = bool(profile.get("referral_discount_active"))
    pct = int(profile.get("referral_discount_percent") or 0)

    invites: list[dict] = []
    try:
        from bot.services.supabase_db import has_db, db_get_referrals_for_inviter

        if has_db():
            invites = db_get_referrals_for_inviter(int(uid), limit=limit)
    except Exception:
        invites = []

    return {
        "ok": True,
        "telegram_id": uid,
        "referrals_count": cnt,
        "referral_discount_active": active,
        "referral_discount_percent": pct,
        "invitees": invites,
    }


@router.post("/api/notify")
async def api_notify(
    req: NotifyRequest,
    ptb=Depends(get_ptb_application),
):
    uid = resolve_telegram_uid(str(req.telegram_id), req.token)
    if not uid:
        raise HTTPException(status_code=401, detail="Foydalanuvchi aniqlanmadi")

    if not req.message or not req.message.strip():
        raise HTTPException(status_code=400, detail="Xabar bo'sh")
    if len(req.message) > NOTIFY_MAX_CHARS:
        raise HTTPException(status_code=400, detail="Xabar 4000 belgi oshmasligi kerak")

    try:
        await ptb.bot.send_message(
            chat_id=int(uid),
            text=req.message,
            parse_mode="HTML",
        )
        try:
            from bot.utils.action_logger import log_action_fire_and_forget

            log_action_fire_and_forget(
                telegram_id=int(uid),
                username=None,
                action_type="NOTIFY",
                details="web:notify",
                metadata={"chars": len(req.message or "")},
            )
        except Exception:
            pass
        return {"ok": True}
    except Exception as e:
        logger.warning("/api/notify failed for %s: %s", uid, e)
        raise HTTPException(status_code=502, detail=f"Telegram xatosi: {str(e)[:200]}")


@router.post("/api/support")
async def api_support(
    req: SupportRequest,
    ptb=Depends(get_ptb_application),
):
    uid = resolve_telegram_uid(str(req.telegram_id), req.token)
    if not uid:
        raise HTTPException(status_code=401, detail="Foydalanuvchi aniqlanmadi")

    msg = (req.message or "").strip()
    if not msg:
        raise HTTPException(status_code=400, detail="Xabar bo'sh")
    if len(msg) > NOTIFY_MAX_CHARS:
        raise HTTPException(status_code=400, detail="Xabar 4000 belgidan oshmasligi kerak")

    admin_ids: list[int] = []
    raw_admin = (os.getenv("ADMIN_USER_ID") or "").strip()
    if raw_admin:
        admin_ids.extend(
            int(x.strip()) for x in raw_admin.split(",") if x.strip().lstrip("-").isdigit()
        )
    support_group = os.getenv("SUPPORT_GROUP_ID", "").strip()
    if support_group and support_group.lstrip("-").isdigit():
        admin_ids.append(int(support_group))

    if not admin_ids:
        admin_ids.append(-1003457224552)

    username = (req.username or "").strip()
    username_text = f"@{username}" if username else "yo'q"
    support_item = None
    try:
        from bot.services.support_service import create_support_request

        support_item = create_support_request(
            user_id=int(uid),
            username=username,
            message=msg,
            source="webapp",
        )
    except Exception as e:
        logger.warning("/api/support save failed: %s", e)

    ticket_line = f"🎫 Request ID: <b>#{support_item['id']}</b>\n" if support_item else ""
    text = (
        "📩 <b>WebApp support so'rovi</b>\n\n"
        f"{ticket_line}"
        f"🆔 User ID: <code>{uid}</code>\n"
        f"👤 Username: {username_text}\n\n"
        f"💬 Xabar:\n{msg}"
    )

    sent = 0
    for chat_id in dict.fromkeys(admin_ids):
        try:
            await ptb.bot.send_message(
                chat_id=chat_id,
                text=text,
                parse_mode="HTML",
            )
            sent += 1
        except Exception as e:
            logger.warning("/api/support forward failed to %s: %s", chat_id, e)

    if sent == 0:
        raise HTTPException(status_code=502, detail="Support xabarini yuborib bo'lmadi")

    return {"ok": True, "forwarded_to": sent}


@router.post("/api/translate")
async def api_translate(req: TranslateRequest):
    if not req.text or not req.text.strip():
        raise HTTPException(status_code=400, detail="Matn bo'sh bo'lishi mumkin emas")
    if len(req.text) > TRANSLATE_MAX_CHARS:
        raise HTTPException(status_code=400, detail="Matn 5000 belgidan oshmasligi kerak")

    valid_dirs = {"uz_en", "en_uz", "ru_uz", "uz_ru", "ru_en", "en_ru"}
    if req.direction not in valid_dirs:
        raise HTTPException(status_code=400, detail=f"Noto'g'ri yo'nalish: {req.direction}")

    uid = _resolve_web_uid_optional(req.telegram_id, req.token)

    quota = None
    consumed = False
    try:
        from bot.services.ai_service import is_meaningfully_changed, translate_text

        if uid:
            from bot.services.plan_limits import CAT_TRANSLATE

            consumed = bool(web_quota_consume_or_raise(int(uid), CAT_TRANSLATE))

        result = await translate_text(req.text, req.direction)
        if not result or result.startswith("Tarjimada xato") or result.startswith("AI model"):
            raise HTTPException(status_code=502, detail=result or "Tarjima bo'sh qaytdi")
        if not is_meaningfully_changed(req.text, result):
            raise HTTPException(status_code=422, detail="Tarjima natijasi original bilan bir xil chiqdi")
        if uid:
            from bot.services.plan_limits import CAT_TRANSLATE

            from bot.services.plan_limits import category_quota_for_response
            from bot.services.user_service import record_service_completion

            record_service_completion(int(uid), CAT_TRANSLATE, "Web translate", skip_quota=bool(consumed))
            quota = category_quota_for_response(int(uid), CAT_TRANSLATE)
            try:
                from bot.utils.action_logger import log_action_fire_and_forget

                log_action_fire_and_forget(
                    telegram_id=int(uid),
                    username=None,
                    action_type="TRANSLATE",
                    details=f"web:{req.direction}",
                    metadata={"direction": req.direction, "chars": len(req.text or "")},
                )
            except Exception:
                pass
        return {"ok": True, "translated_text": result, "direction": req.direction, "quota": quota}
    except HTTPException:
        raise
    except Exception as e:
        logger.error("Translate API error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=f"Tarjima serveri xatosi: {str(e)[:200]}")


@router.post("/api/translate_auto")
async def api_translate_auto(req: TranslateAutoRequest):
    """
    Auto source detection (no source selection).
    Client provides only target_lang (uz|ru|en).
    """
    if not req.text or not req.text.strip():
        raise HTTPException(status_code=400, detail="Matn bo'sh bo'lishi mumkin emas")
    if len(req.text) > TRANSLATE_MAX_CHARS:
        raise HTTPException(status_code=400, detail="Matn 5000 belgidan oshmasligi kerak")

    tgt = (req.target_lang or "").strip().lower()
    if tgt not in {"uz", "ru", "en"}:
        raise HTTPException(status_code=400, detail="target_lang noto'g'ri (uz|ru|en)")

    # Cache key (text+target) to avoid repeated AI calls
    cache_key = None
    try:
        if TRANSLATE_AUTO_CACHE_TTL_SECONDS > 0:
            h = hashlib.sha1((req.text or "").encode("utf-8", errors="ignore")).hexdigest()  # nosec - non-crypto use
            cache_key = f"tr:auto:{h}:{tgt}"
            cached = await redis_cache_get_json(cache_key)
            if cached and isinstance(cached, dict) and cached.get("ok") and cached.get("translated_text"):
                return cached
    except Exception:
        cache_key = None

    # best-effort source detection
    src = "uz"
    try:
        from langdetect import detect  # type: ignore

        sample = (req.text or "").strip().replace("\u0000", "")[:2500]
        d = (detect(sample) or "").lower()
        if d.startswith("ru"):
            src = "ru"
        elif d.startswith("en"):
            src = "en"
        else:
            src = "uz"
    except Exception:
        src = "uz"

    uid = _resolve_web_uid_optional(req.telegram_id, req.token)
    quota = None
    consumed = False
    try:
        if uid:
            from bot.services.plan_limits import CAT_TRANSLATE

            consumed = bool(web_quota_consume_or_raise(int(uid), CAT_TRANSLATE))

        if src == tgt:
            result = req.text.strip()
            direction = "none"
        else:
            direction_map = {
                ("uz", "en"): "uz_en",
                ("en", "uz"): "en_uz",
                ("ru", "uz"): "ru_uz",
                ("uz", "ru"): "uz_ru",
                ("ru", "en"): "ru_en",
                ("en", "ru"): "en_ru",
            }
            direction = direction_map.get((src, tgt))
            if not direction:
                raise HTTPException(status_code=400, detail=f"translate yo'nalishi topilmadi: {src}->{tgt}")
            from bot.services.ai_service import is_meaningfully_changed, translate_text

            result = await translate_text(req.text, direction)
            if not result or result.startswith("Tarjimada xato") or result.startswith("Tarjima vaqtincha") or result.startswith("AI model"):
                raise HTTPException(status_code=502, detail=result or "Tarjima bo'sh qaytdi")
            if not is_meaningfully_changed(req.text, result):
                raise HTTPException(status_code=422, detail="Tarjima natijasi original bilan bir xil chiqdi")

        if uid and direction != "none":
            from bot.services.plan_limits import CAT_TRANSLATE

            from bot.services.plan_limits import category_quota_for_response
            from bot.services.user_service import record_service_completion

            record_service_completion(int(uid), CAT_TRANSLATE, "Web translate auto", skip_quota=bool(consumed))
            quota = category_quota_for_response(int(uid), CAT_TRANSLATE)
            try:
                from bot.utils.action_logger import log_action_fire_and_forget

                log_action_fire_and_forget(
                    telegram_id=int(uid),
                    username=None,
                    action_type="TRANSLATE_AUTO",
                    details=f"web:{src}->{tgt}",
                    metadata={"source_lang": src, "target_lang": tgt, "chars": len(req.text or "")},
                )
            except Exception:
                pass
        resp = {
            "ok": True,
            "translated_text": result,
            "source_lang": src,
            "target_lang": tgt,
            "direction": direction,
            "quota": quota,
        }
        try:
            if cache_key:
                await redis_cache_set_json(cache_key, resp, TRANSLATE_AUTO_CACHE_TTL_SECONDS)
        except Exception:
            pass
        return resp
    except HTTPException:
        raise
    except Exception as e:
        logger.error("Translate auto API error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=f"Tarjima serveri xatosi: {str(e)[:200]}")



@router.post("/api/spellcheck")
async def api_spellcheck(req: SpellcheckRequest):
    if not req.text or not req.text.strip():
        raise HTTPException(status_code=400, detail="Matn bo'sh bo'lishi mumkin emas")
    if len(req.text) > SPELLCHECK_MAX_CHARS:
        raise HTTPException(
            status_code=400,
            detail=f"Matn {SPELLCHECK_MAX_CHARS} belgidan oshmasligi kerak",
        )

    uid = _resolve_web_uid_optional(req.telegram_id, req.token)

    quota = None
    try:
        from bot.services.ai_service import check_spelling_text, count_words_text, is_meaningfully_changed

        src = req.text.strip()
        wc = count_words_text(src)
        key = spellcheck_cache_key(src)
        consumed = False
        if uid:
            from bot.services.plan_limits import CAT_SPELL

            consumed = bool(web_quota_consume_or_raise(int(uid), CAT_SPELL))
        hit = spellcheck_cache_get(key)
        if hit is not None:
            corrected, fixes = hit
            fc = int(fixes or 0)
            if uid:
                from bot.services.plan_limits import CAT_SPELL

                from bot.services.plan_limits import category_quota_for_response
                from bot.services.user_service import record_service_completion

                record_service_completion(int(uid), CAT_SPELL, "Web spell text", skip_quota=bool(consumed))
                quota = category_quota_for_response(int(uid), CAT_SPELL)
                try:
                    from bot.utils.action_logger import log_action_fire_and_forget

                    log_action_fire_and_forget(
                        telegram_id=int(uid),
                        username=None,
                        action_type="SPELL",
                        details="web:text(cache_hit)",
                        metadata={"chars": len(src or ""), "word_count": wc, "fixed": fc},
                    )
                except Exception:
                    pass
            return {
                "ok": True,
                "corrected_text": corrected,
                "fixed": fc,
                "error_count": fc,
                "word_count": wc,
                "quota": quota,
            }

        corrected, fixes = await check_spelling_text(src)
        if corrected is None:
            raise HTTPException(status_code=502, detail="Natija bo'sh qaytdi")
        if int(fixes or 0) == 0 and is_meaningfully_changed(src, corrected):
            fixes = 1

        spellcheck_cache_set(key, corrected, int(fixes or 0))
        fc = int(fixes or 0)
        if uid:
            from bot.services.plan_limits import CAT_SPELL

            from bot.services.plan_limits import category_quota_for_response
            from bot.services.user_service import record_service_completion

            record_service_completion(int(uid), CAT_SPELL, "Web spell text", skip_quota=bool(consumed))
            quota = category_quota_for_response(int(uid), CAT_SPELL)
            try:
                from bot.utils.action_logger import log_action_fire_and_forget

                log_action_fire_and_forget(
                    telegram_id=int(uid),
                    username=None,
                    action_type="SPELL",
                    details="web:text",
                    metadata={"chars": len(src or ""), "word_count": wc, "fixed": fc},
                )
            except Exception:
                pass
        return {
            "ok": True,
            "corrected_text": corrected,
            "fixed": fc,
            "error_count": fc,
            "word_count": wc,
            "quota": quota,
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error("Spellcheck API error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=f"Imlo serveri xatosi: {str(e)[:200]}")


@router.post("/api/spellcheck_file")
async def api_spellcheck_file(
    file: UploadFile = File(...),
    token: Optional[str] = Form(None),
    telegram_id: Optional[str] = Form(None),
    notify_telegram: str = Form("false"),
    ptb=Depends(get_ptb_application),
):
    """
    Extract text from .txt / .docx / .pptx / .pdf → spell-check → JSON.
    Optional: send corrected UTF-8 .txt to the user's Telegram chat.
    """
    uid = resolve_telegram_uid(telegram_id, token)
    consumed = False
    if uid:
        # Enforce premium/limits BEFORE heavy AI work.
        from bot.services.plan_limits import CAT_SPELL

        consumed = bool(web_quota_consume_or_raise(int(uid), CAT_SPELL))
    max_b = get_settings().max_upload_bytes
    raw = await file.read()
    do_notify = str(notify_telegram).lower() in ("1", "true", "yes", "on")
    if uid:
        do_notify = True
    uid_int = int(uid) if uid else None
    logger.info(
        "POST /api/spellcheck_file name=%s bytes=%s notify=%s uid=%s",
        file.filename,
        len(raw),
        do_notify,
        uid or "-",
    )
    if not raw:
        raise HTTPException(status_code=400, detail="Bo'sh fayl")
    if len(raw) > max_b:
        raise HTTPException(status_code=413, detail="Fayl juda katta")

    # Filename fallback by content-type (ba'zi klientlar filename ni noto'g'ri yuboradi)
    fname = file.filename or "upload.txt"
    if "." not in fname:
        ct = (file.content_type or "").lower()
        if "wordprocessingml" in ct:
            fname = fname + ".docx"
        elif "presentationml" in ct or "powerpoint" in ct:
            fname = fname + ".pptx"
        elif "pdf" in ct:
            fname = fname + ".pdf"
        elif "text" in ct:
            fname = fname + ".txt"

    try:
        from bot.services.document_text_extract import extract_plain_text_from_bytes

        text = extract_plain_text_from_bytes(fname, raw)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e

    if not (text or "").strip():
        raise HTTPException(status_code=422, detail="Fayldan matn ajratilmadi")

    src = text.strip()
    if len(src) > SPELLCHECK_FILE_MAX_CHARS:
        raise HTTPException(
            status_code=400,
            detail=f"Fayldan olingan matn {SPELLCHECK_FILE_MAX_CHARS} belgidan oshmasligi kerak",
        )

    from bot.services.ai_service import check_spelling_text, count_words_text, is_meaningfully_changed

    wc = count_words_text(src)
    key = spellcheck_cache_key(src)
    hit = spellcheck_cache_get(key)
    if hit is not None:
        corrected, fixes = hit
    else:
        corrected, fixes = await check_spelling_text(src)
        if corrected is None:
            raise HTTPException(status_code=502, detail="Natija bo'sh qaytdi")
        if int(fixes or 0) == 0 and is_meaningfully_changed(src, corrected):
            fixes = 1
        spellcheck_cache_set(key, corrected, int(fixes or 0))

    fc = int(fixes or 0)
    quota = None
    if uid_int:
        try:
            from bot.services.plan_limits import CAT_SPELL, category_quota_for_response
            from bot.services.user_service import record_service_completion

            record_service_completion(int(uid_int), CAT_SPELL, "Web spell file", skip_quota=bool(consumed))
            quota = category_quota_for_response(int(uid_int), CAT_SPELL)
        except Exception:
            quota = None
        try:
            from bot.utils.action_logger import log_action_fire_and_forget

            log_action_fire_and_forget(
                telegram_id=int(uid_int),
                username=None,
                action_type="SPELL",
                details=f"web:file:{os.path.splitext(fname)[1].lower() or 'unknown'}",
                metadata={"filename": fname, "bytes": len(raw), "word_count": wc, "fixed": fc},
            )
        except Exception:
            pass
    out = {
        "ok": True,
        "corrected_text": corrected,
        "fixed": fc,
        "error_count": fc,
        "word_count": wc,
        "filename": file.filename or "document",
        "quota": quota,
    }

    if do_notify and uid:
        try:
            from bot.services.user_service import get_chat_id
            from bot.services.ai_service import check_spelling_gemini, check_spelling_pptx

            chat_id = get_chat_id(int(uid)) or int(uid)
            # Use file-based fixed-count when possible (docx/pptx pipelines return better numbers).
            fixed_for_summary = int(fc or 0)
            base_name, ext = os.path.splitext(fname)
            ext = ext.lower() if ext else ".txt"
            send_name = f"{base_name}_imlo_tuzatilgan{ext}"
            data_buf: io.BytesIO | None = None
            tmp_in = None
            tmp_out = None

            if ext == ".docx":
                fd_in, tmp_in = tempfile.mkstemp(suffix=".docx", prefix="web_spell_")
                os.close(fd_in)
                with open(tmp_in, "wb") as wf:
                    wf.write(raw)
                tmp_out, fixed2, _ = await check_spelling_gemini(tmp_in)
                try:
                    fixed_for_summary = max(int(fixed_for_summary), int(fixed2 or 0))
                except Exception:
                    pass
                if not tmp_out or not os.path.exists(tmp_out):
                    # Fallback to plain text file if docx build failed.
                    send_name = f"{base_name}_imlo_tuzatilgan.txt"
                    data_buf = io.BytesIO((corrected or "").encode("utf-8"))
            elif ext == ".pptx":
                fd_in, tmp_in = tempfile.mkstemp(suffix=".pptx", prefix="web_spell_")
                os.close(fd_in)
                with open(tmp_in, "wb") as wf:
                    wf.write(raw)
                tmp_out, fixed2, _ = await check_spelling_pptx(tmp_in)
                try:
                    fixed_for_summary = max(int(fixed_for_summary), int(fixed2 or 0))
                except Exception:
                    pass
                if not tmp_out or not os.path.exists(tmp_out):
                    send_name = f"{base_name}_imlo_tuzatilgan.txt"
                    data_buf = io.BytesIO((corrected or "").encode("utf-8"))
            elif ext == ".pdf":
                fd_out, tmp_out = tempfile.mkstemp(suffix=".pdf", prefix="web_spell_")
                os.close(fd_out)
                _build_pdf_from_text(corrected or "", tmp_out)
            else:
                send_name = f"{base_name}_imlo_tuzatilgan.txt"
                data_buf = io.BytesIO((corrected or "").encode("utf-8"))

            summary = f"✅ Imlo tekshiruvi (WebApp)\nSo'zlar: {wc}\nTopilgan xatolar: {fixed_for_summary}"
            await ptb.bot.send_message(chat_id=chat_id, text=summary)

            if tmp_out and os.path.exists(tmp_out):
                with open(tmp_out, "rb") as fp:
                    await ptb.bot.send_document(
                        chat_id=chat_id,
                        document=InputFile(fp, filename=send_name),
                        caption=f"Tuzatilgan fayl ({ext[1:].upper() if ext.startswith('.') else ext.upper()})",
                    )
            else:
                buf = data_buf or io.BytesIO((corrected or "").encode("utf-8"))
                buf.seek(0)
                await ptb.bot.send_document(
                    chat_id=chat_id,
                    document=InputFile(buf, filename=send_name),
                    caption="Tuzatilgan matn (TXT)",
                )
            for p in (tmp_in, tmp_out):
                if p and os.path.exists(p):
                    try:
                        os.remove(p)
                    except Exception:
                        pass
        except Exception as e:
            logger.warning("spellcheck_file Telegram: %s", e)

    return out


@router.post("/api/spellcheck_file_download")
async def api_spellcheck_file_download(
    file: UploadFile = File(...),
    token: Optional[str] = Form(None),
    telegram_id: Optional[str] = Form(None),
):
    """
    Download corrected file with ORIGINAL extension:
      - docx -> *_tuzatilgan.docx  (best-effort: AI rebuild if available)
      - pptx -> *_tuzatilgan.pptx  (best-effort)
      - pdf  -> *_tuzatilgan.pdf   (generated from corrected text)
      - xlsx -> *_tuzatilgan.xlsx  (generated sheet with corrected text lines)
      - txt  -> *_tuzatilgan.txt
    """
    uid = resolve_telegram_uid(telegram_id, token)
    consumed = False
    if uid:
        # Enforce premium/limits BEFORE heavy AI work.
        from bot.services.plan_limits import CAT_SPELL

        consumed = bool(web_quota_consume_or_raise(int(uid), CAT_SPELL))
    max_b = get_settings().max_upload_bytes
    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Bo'sh fayl")
    if len(raw) > max_b:
        raise HTTPException(status_code=413, detail="Fayl juda katta")

    fname = file.filename or "upload.txt"
    if "." not in fname:
        ct = (file.content_type or "").lower()
        if "spreadsheetml" in ct or "excel" in ct:
            fname = fname + ".xlsx"
        elif "wordprocessingml" in ct:
            fname = fname + ".docx"
        elif "presentationml" in ct or "powerpoint" in ct:
            fname = fname + ".pptx"
        elif "pdf" in ct:
            fname = fname + ".pdf"
        else:
            fname = fname + ".txt"

    try:
        from bot.services.document_text_extract import extract_plain_text_from_bytes

        text = extract_plain_text_from_bytes(fname, raw)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e

    if not (text or "").strip():
        raise HTTPException(status_code=422, detail="Fayldan matn ajratilmadi")

    src = text.strip()
    if len(src) > SPELLCHECK_FILE_MAX_CHARS:
        raise HTTPException(
            status_code=400,
            detail=f"Fayldan olingan matn {SPELLCHECK_FILE_MAX_CHARS} belgidan oshmasligi kerak",
        )

    from bot.services.ai_service import check_spelling_text, is_meaningfully_changed

    corrected, fixes = await check_spelling_text(src)
    if corrected is None:
        raise HTTPException(status_code=502, detail="Natija bo'sh qaytdi")
    if int(fixes or 0) == 0 and is_meaningfully_changed(src, corrected):
        fixes = 1

    base_name, ext = os.path.splitext(fname)
    ext = (ext or ".txt").lower()
    safe_base = os.path.basename(base_name).strip() or "document"
    out_name = f"{safe_base}_tuzatilgan{ext}"

    data_buf: io.BytesIO | None = None
    tmp_in = None
    tmp_out = None
    content_type = "application/octet-stream"

    try:
        if ext == ".docx":
            # Prefer AI rebuild (keeps more formatting than plain text)
            try:
                from bot.services.ai_service import check_spelling_gemini

                fd_in, tmp_in = tempfile.mkstemp(suffix=".docx", prefix="web_spell_dl_")
                os.close(fd_in)
                with open(tmp_in, "wb") as wf:
                    wf.write(raw)
                tmp_out, _, _ = await check_spelling_gemini(tmp_in)
            except Exception:
                tmp_out = None
            if tmp_out and os.path.exists(tmp_out):
                content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                with open(tmp_out, "rb") as fp:
                    data_buf = io.BytesIO(fp.read())
            else:
                # Fallback: build a simple docx from corrected text
                try:
                    from docx import Document

                    doc = Document()
                    for line in (corrected or "").splitlines():
                        doc.add_paragraph(line)
                    bio = io.BytesIO()
                    doc.save(bio)
                    bio.seek(0)
                    data_buf = bio
                    content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                except Exception:
                    out_name = f"{safe_base}_tuzatilgan.txt"
                    data_buf = io.BytesIO((corrected or "").encode("utf-8"))
                    content_type = "text/plain; charset=utf-8"

        elif ext == ".pptx":
            try:
                from bot.services.ai_service import check_spelling_pptx

                fd_in, tmp_in = tempfile.mkstemp(suffix=".pptx", prefix="web_spell_dl_")
                os.close(fd_in)
                with open(tmp_in, "wb") as wf:
                    wf.write(raw)
                tmp_out, _, _ = await check_spelling_pptx(tmp_in)
            except Exception:
                tmp_out = None
            if tmp_out and os.path.exists(tmp_out):
                content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                with open(tmp_out, "rb") as fp:
                    data_buf = io.BytesIO(fp.read())
            else:
                # Fallback: create a simple pptx with one slide text box
                try:
                    from pptx import Presentation
                    from pptx.util import Inches

                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(5.0))
                    tf = tx.text_frame
                    for i, line in enumerate((corrected or "").splitlines()[:200]):
                        if i == 0:
                            tf.text = line
                        else:
                            tf.add_paragraph().text = line
                    bio = io.BytesIO()
                    prs.save(bio)
                    bio.seek(0)
                    data_buf = bio
                    content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                except Exception:
                    out_name = f"{safe_base}_tuzatilgan.txt"
                    data_buf = io.BytesIO((corrected or "").encode("utf-8"))
                    content_type = "text/plain; charset=utf-8"

        elif ext == ".pdf":
            fd_out, tmp_out = tempfile.mkstemp(suffix=".pdf", prefix="web_spell_dl_")
            os.close(fd_out)
            _build_pdf_from_text(corrected or "", tmp_out)
            content_type = "application/pdf"
            with open(tmp_out, "rb") as fp:
                data_buf = io.BytesIO(fp.read())

        elif ext == ".xlsx":
            try:
                from openpyxl import Workbook

                wb = Workbook()
                ws = wb.active
                ws.title = "Tuzatilgan"
                for i, line in enumerate((corrected or "").splitlines(), start=1):
                    ws.cell(row=i, column=1, value=line)
                    if i >= 50000:
                        break
                bio = io.BytesIO()
                wb.save(bio)
                bio.seek(0)
                data_buf = bio
                content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            except Exception:
                out_name = f"{safe_base}_tuzatilgan.txt"
                data_buf = io.BytesIO((corrected or "").encode("utf-8"))
                content_type = "text/plain; charset=utf-8"

        else:
            out_name = f"{safe_base}_tuzatilgan.txt"
            data_buf = io.BytesIO((corrected or "").encode("utf-8"))
            content_type = "text/plain; charset=utf-8"

        if uid:
            # Finalize usage (already pre-consumed for non-admins).
            try:
                from bot.services.plan_limits import CAT_SPELL, category_quota_for_response
                from bot.services.user_service import record_service_completion

                record_service_completion(int(uid), CAT_SPELL, "Web spell file download", skip_quota=bool(consumed))
                _ = category_quota_for_response(int(uid), CAT_SPELL)
            except Exception:
                pass

        data_buf.seek(0)
        return StreamingResponse(
            data_buf,
            media_type=content_type,
            headers={"Content-Disposition": f'attachment; filename="{out_name}"'},
        )
    finally:
        for p in (tmp_in, tmp_out):
            if p and os.path.exists(p):
                try:
                    os.remove(p)
                except Exception:
                    pass


@router.post("/api/translit_file_download")
async def api_translit_file_download(
    file: UploadFile = File(...),
    token: Optional[str] = Form(None),
    telegram_id: Optional[str] = Form(None),
    notify_telegram: str = Form("false"),
    ptb=Depends(get_ptb_application),
):
    """
    Download auto Cyrillic↔Latin converted file with ORIGINAL extension.
    If detection is unclear, returns original as-is.

    Supported:
      - docx -> *_ogirilgan.docx
      - pptx -> *_ogirilgan.pptx
      - pdf  -> *_ogirilgan.pdf   (generated from converted text; fallback to original if unknown)
      - xlsx -> *_ogirilgan.xlsx
      - txt  -> *_ogirilgan.txt
    """
    uid = resolve_telegram_uid(telegram_id, token)
    consumed = False
    do_notify = str(notify_telegram).lower() in ("1", "true", "yes", "on")
    if uid:
        do_notify = True
    if uid:
        # Enforce premium/limits BEFORE heavy file rebuild.
        from bot.services.plan_limits import CAT_TRANSLIT

        consumed = bool(web_quota_consume_or_raise(int(uid), CAT_TRANSLIT))
    max_b = get_settings().max_upload_bytes
    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Bo'sh fayl")
    if len(raw) > max_b:
        raise HTTPException(status_code=413, detail="Fayl juda katta")

    fname = file.filename or "upload.txt"
    if "." not in fname:
        ct = (file.content_type or "").lower()
        if "spreadsheetml" in ct or "excel" in ct:
            fname = fname + ".xlsx"
        elif "wordprocessingml" in ct:
            fname = fname + ".docx"
        elif "presentationml" in ct or "powerpoint" in ct:
            fname = fname + ".pptx"
        elif "pdf" in ct:
            fname = fname + ".pdf"
        else:
            fname = fname + ".txt"

    base_name, ext = os.path.splitext(fname)
    ext = (ext or ".txt").lower()
    safe_base = os.path.basename(base_name).strip() or "document"
    out_name = f"{safe_base}_ogirilgan{ext}"

    try:
        from bot.services.document_text_extract import extract_plain_text_from_bytes

        text = extract_plain_text_from_bytes(fname, raw)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Fayl o'qilmadi: {str(e)[:200]}") from e

    auto = auto_cyrillic_latin(text or "")
    if auto.direction == "none":
        # return original unchanged (but with our suffix name)
        return StreamingResponse(
            io.BytesIO(raw),
            media_type=file.content_type or "application/octet-stream",
            headers={"Content-Disposition": f'attachment; filename="{out_name}"'},
        )

    from bot.services.transliterate_service import transliterate

    data_buf: io.BytesIO | None = None
    tmp_out = None
    content_type = "application/octet-stream"

    try:
        if ext == ".docx":
            try:
                from docx import Document

                doc = Document(io.BytesIO(raw))
                for para in doc.paragraphs:
                    for run in para.runs:
                        if run.text:
                            run.text = transliterate(run.text, auto.direction)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    if run.text:
                                        run.text = transliterate(run.text, auto.direction)
                bio = io.BytesIO()
                doc.save(bio)
                bio.seek(0)
                data_buf = bio
                content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            except Exception:
                out_name = f"{safe_base}_ogirilgan.txt"
                data_buf = io.BytesIO((auto.result or "").encode("utf-8"))
                content_type = "text/plain; charset=utf-8"

        elif ext == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(io.BytesIO(raw))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.text:
                                        run.text = transliterate(run.text, auto.direction)
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    for para in cell.text_frame.paragraphs:
                                        for run in para.runs:
                                            if run.text:
                                                run.text = transliterate(run.text, auto.direction)
                bio = io.BytesIO()
                prs.save(bio)
                bio.seek(0)
                data_buf = bio
                content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            except Exception:
                out_name = f"{safe_base}_ogirilgan.txt"
                data_buf = io.BytesIO((auto.result or "").encode("utf-8"))
                content_type = "text/plain; charset=utf-8"

        elif ext == ".xlsx":
            try:
                from openpyxl import load_workbook

                wb = load_workbook(io.BytesIO(raw))
                for ws in wb.worksheets:
                    for row in ws.iter_rows():
                        for cell in row:
                            v = cell.value
                            if isinstance(v, str) and v.strip():
                                cell.value = transliterate(v, auto.direction)
                bio = io.BytesIO()
                wb.save(bio)
                bio.seek(0)
                data_buf = bio
                content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            except Exception:
                out_name = f"{safe_base}_ogirilgan.txt"
                data_buf = io.BytesIO((auto.result or "").encode("utf-8"))
                content_type = "text/plain; charset=utf-8"

        elif ext == ".pdf":
            # Rebuild PDF from converted plain text (formatting not preserved)
            fd_out, tmp_out = tempfile.mkstemp(suffix=".pdf", prefix="web_translit_dl_")
            os.close(fd_out)
            _build_pdf_from_text(auto.result or "", tmp_out)
            content_type = "application/pdf"
            with open(tmp_out, "rb") as fp:
                data_buf = io.BytesIO(fp.read())

        else:
            out_name = f"{safe_base}_ogirilgan.txt"
            data_buf = io.BytesIO((auto.result or "").encode("utf-8"))
            content_type = "text/plain; charset=utf-8"

        if uid:
            # Finalize usage (already pre-consumed for non-admins).
            try:
                from bot.services.plan_limits import CAT_TRANSLIT, category_quota_for_response
                from bot.services.user_service import record_service_completion

                record_service_completion(int(uid), CAT_TRANSLIT, "Web translit file download", skip_quota=bool(consumed))
                _ = category_quota_for_response(int(uid), CAT_TRANSLIT)
            except Exception:
                pass
            try:
                from bot.utils.action_logger import log_action_fire_and_forget

                log_action_fire_and_forget(
                    telegram_id=int(uid),
                    username=None,
                    action_type="TRANSLIT_AUTO",
                    details=f"web:file:{ext}",
                    metadata={
                        "filename": fname,
                        "bytes": len(raw),
                        "detected": auto.detected,
                        "direction": auto.direction,
                        "letters": (auto.stats or {}).get("letters"),
                    },
                )
            except Exception:
                pass

        (data_buf or io.BytesIO(b"")).seek(0)
        buf = data_buf or io.BytesIO(b"")
        buf.seek(0)

        if do_notify and uid:
            try:
                from telegram import InputFile

                chat_id = int(uid)
                caption = "✅ Krill↔Lotin tayyor (WebApp)"
                # send copy to Telegram so user can continue flows in bot
                await ptb.bot.send_document(
                    chat_id=chat_id,
                    document=InputFile(io.BytesIO(buf.getvalue()), filename=out_name),
                    caption=caption,
                )
            except Exception:
                pass

        return StreamingResponse(
            io.BytesIO(buf.getvalue()),
            media_type=content_type,
            headers={"Content-Disposition": f'attachment; filename="{out_name}"'},
        )
    finally:
        if tmp_out and os.path.exists(tmp_out):
            try:
                os.remove(tmp_out)
            except Exception:
                pass


@router.post("/api/process_file")
async def api_process_file(
    file: UploadFile = File(...),
    action: str = Form(..., description="spellcheck|translit_auto|translate_auto"),
    target_lang: Optional[str] = Form(None, description="uz|ru|en (for translate_auto)"),
    token: Optional[str] = Form(None),
    telegram_id: Optional[str] = Form(None),
    notify_telegram: str = Form("false"),
    ptb=Depends(get_ptb_application),
):
    """
    Unified pipeline:
      upload -> (Supabase Storage if configured) -> extract text -> process -> structured JSON
    """
    uid = resolve_telegram_uid(telegram_id, token)
    uid_int = int(uid) if uid else None
    do_notify = str(notify_telegram).lower() in ("1", "true", "yes", "on")
    if uid_int:
        do_notify = True

    try:
        raw = await read_upload_limited(file)
    except EmptyUploadError:
        raise HTTPException(status_code=400, detail="Bo'sh fayl")
    except UploadTooLargeError as e:
        raise HTTPException(status_code=413, detail=str(e))

    fname = file.filename or "upload.txt"
    # content-type based extension fallback (same approach as spellcheck_file)
    if "." not in fname:
        ct = (file.content_type or "").lower()
        if "spreadsheetml" in ct or "excel" in ct:
            fname = fname + ".xlsx"
        elif "wordprocessingml" in ct:
            fname = fname + ".docx"
        elif "presentationml" in ct or "powerpoint" in ct:
            fname = fname + ".pptx"
        elif "pdf" in ct:
            fname = fname + ".pdf"
        elif "text" in ct:
            fname = fname + ".txt"

    st = None
    try:
        if uid_int:
            st = upload_bytes_for_user(telegram_id=uid_int, filename=fname, raw=raw)
    except Exception:
        st = None

    try:
        from bot.services.document_text_extract import extract_plain_text_from_bytes

        text = extract_plain_text_from_bytes(fname, raw)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e

    if not (text or "").strip():
        raise HTTPException(status_code=422, detail="Fayldan matn ajratilmadi")

    act = (action or "").strip().lower()
    if act not in {"spellcheck", "translit_auto", "translate_auto"}:
        raise HTTPException(status_code=400, detail="action noto'g'ri")

    result_text = ""
    meta: dict = {"filename": fname, "chars": len(text)}
    quota = None
    consumed = False
    cat = None

    if act == "spellcheck":
        if uid_int:
            from bot.services.plan_limits import CAT_SPELL

            cat = CAT_SPELL
            consumed = bool(web_quota_consume_or_raise(int(uid_int), cat))
        try:
            from bot.services.ai_service import check_spelling_text

            corrected, fixes = await check_spelling_text(text)
            result_text = corrected
            meta["fixed"] = int(fixes or 0)
        except HTTPException:
            raise
        except Exception as e:
            logger.error("process_file spellcheck error: %s", e, exc_info=True)
            raise HTTPException(status_code=502, detail=f"Imlo tekshirish xatosi: {str(e)[:200]}")

    elif act == "translit_auto":
        try:
            res = auto_cyrillic_latin(text)
            result_text = res.result
            meta["detected"] = res.detected
            meta["direction"] = res.direction
            meta["stats"] = res.stats
            if uid_int and res.direction != "none":
                from bot.services.plan_limits import CAT_TRANSLIT

                cat = CAT_TRANSLIT
                consumed = bool(web_quota_consume_or_raise(int(uid_int), cat))
        except HTTPException:
            raise
        except Exception as e:
            logger.error("process_file translit_auto error: %s", e, exc_info=True)
            raise HTTPException(status_code=500, detail=f"Krill↔Lotin xatosi: {str(e)[:200]}")

    else:
        # translate_auto: detect source language, ask only for target_lang if needed
        tgt = (target_lang or "").strip().lower()
        if tgt not in {"uz", "ru", "en"}:
            raise HTTPException(status_code=400, detail="target_lang kerak (uz|ru|en)")

        # Detect source (best-effort)
        src = "uz"
        try:
            from langdetect import detect  # type: ignore

            sample = (text or "").strip().replace("\u0000", "")
            sample = sample[:2500]
            d = (detect(sample) or "").lower()
            if d.startswith("ru"):
                src = "ru"
            elif d.startswith("en"):
                src = "en"
            else:
                src = "uz"
        except Exception:
            src = "uz"

        if src == tgt:
            result_text = text
            meta["source_lang"] = src
            meta["target_lang"] = tgt
            meta["direction"] = "none"
        else:
            direction_map = {
                ("uz", "en"): "uz_en",
                ("en", "uz"): "en_uz",
                ("ru", "uz"): "ru_uz",
                ("uz", "ru"): "uz_ru",
                ("ru", "en"): "ru_en",
                ("en", "ru"): "en_ru",
            }
            direction = direction_map.get((src, tgt))
            if not direction:
                raise HTTPException(status_code=400, detail=f"translate yo'nalishi topilmadi: {src}->{tgt}")
            try:
                if uid_int:
                    # Enforce premium/limits BEFORE heavy AI work.
                    from bot.services.plan_limits import CAT_TRANSLATE

                    cat = CAT_TRANSLATE
                    consumed = bool(web_quota_consume_or_raise(int(uid_int), cat))

                from bot.services.ai_service import translate_text

                out = await translate_text(text, direction)
                if not out or out.startswith("Tarjima vaqtincha") or out.startswith("Tarjimada xato"):
                    raise HTTPException(status_code=502, detail=out or "Tarjima bo'sh qaytdi")
                result_text = out
                meta["source_lang"] = src
                meta["target_lang"] = tgt
                meta["direction"] = direction
            except HTTPException:
                raise
            except Exception as e:
                logger.error(
                    "process_file translate_auto error: %s (src=%s tgt=%s dir=%s fname=%s bytes=%s chars=%s)",
                    e,
                    src,
                    tgt,
                    direction,
                    fname,
                    len(raw),
                    len(text or ""),
                    exc_info=True,
                )
                raise HTTPException(status_code=502, detail=f"Tarjima xatosi: {str(e)[:200]}")

    signed = None
    if st and st.bucket and st.path:
        try:
            signed = create_signed_url(bucket=st.bucket, path=st.path, expires_in=3600)
        except Exception:
            signed = None

    # Quota (single source via plan_limits / Supabase buckets if configured)
    if uid_int and cat:
        try:
            from bot.services.plan_limits import category_quota_for_response
            from bot.services.user_service import record_service_completion

            record_service_completion(int(uid_int), cat, f"Web file:{act}", skip_quota=bool(consumed))
            quota = category_quota_for_response(int(uid_int), cat)
        except HTTPException:
            # This means user is blocked and we should not return success.
            raise
        except Exception:
            # Never fail the main response if quota snapshot couldn't be built.
            quota = None

    # Telemetry (never block response)
    try:
        if uid_int:
            from bot.utils.action_logger import log_action_fire_and_forget

            log_action_fire_and_forget(
                telegram_id=int(uid_int),
                username=None,
                action_type="PROCESS_FILE",
                details=f"web:{act}",
                metadata={
                    "filename": fname,
                    "bytes": len(raw),
                    "chars": len(text or ""),
                    "target_lang": target_lang,
                    "storage_bucket": st.bucket if st else None,
                },
            )
    except Exception:
        pass

    # Optional: send result to Telegram (so user can continue in bot flows)
    if do_notify and uid_int:
        async def _notify_bg() -> None:
            try:
                from telegram import InputFile

                chat_id = int(uid_int)
                base_name = os.path.splitext(os.path.basename(fname))[0] or "document"
                ext = os.path.splitext(fname)[1].lower() or ".txt"
                out_name = f"{base_name}_{act}{ext}"
                cap = (
                    "✅ Tayyor (WebApp)\n"
                    + (
                        "🌍 Tarjima"
                        if act == "translate_auto"
                        else "🔤 Krill↔Lotin"
                        if act == "translit_auto"
                        else "✏️ Imlo"
                    )
                )
                buf: io.BytesIO | None = None

                if act == "translate_auto":
                    # Try to keep original extension (best-effort). This may be slow — run in background.
                    try:
                        tgt = (target_lang or "").strip().lower()
                        tgt = tgt if tgt in ("uz", "ru", "en") else "uz"
                        direction = (meta or {}).get("direction") if isinstance(meta, dict) else None
                        direction = str(direction or "").strip()
                        if ext == ".docx":
                            from bot.services.ai_service import translate_document_gemini

                            fd_in, tmp_in = tempfile.mkstemp(suffix=".docx", prefix="web_trl_notify_")
                            os.close(fd_in)
                            with open(tmp_in, "wb") as wf:
                                wf.write(raw)
                            tmp_out = await translate_document_gemini(tmp_in, tgt)
                            if tmp_out and os.path.exists(tmp_out):
                                with open(tmp_out, "rb") as fp:
                                    buf = io.BytesIO(fp.read())
                            for p in (tmp_in, tmp_out):
                                try:
                                    if p and os.path.exists(p):
                                        os.remove(p)
                                except Exception:
                                    pass
                        elif ext == ".pptx":
                            from bot.services.ai_service import translate_pptx

                            if not direction:
                                raise Exception("direction yo'q")
                            fd_in, tmp_in = tempfile.mkstemp(suffix=".pptx", prefix="web_trl_notify_")
                            os.close(fd_in)
                            with open(tmp_in, "wb") as wf:
                                wf.write(raw)
                            tmp_out = await translate_pptx(tmp_in, direction, tgt)
                            if tmp_out and os.path.exists(tmp_out):
                                with open(tmp_out, "rb") as fp:
                                    buf = io.BytesIO(fp.read())
                            for p in (tmp_in, tmp_out):
                                try:
                                    if p and os.path.exists(p):
                                        os.remove(p)
                                except Exception:
                                    pass
                        elif ext == ".pdf":
                            fd_out, tmp_out = tempfile.mkstemp(suffix=".pdf", prefix="web_trl_notify_")
                            os.close(fd_out)
                            _build_pdf_from_text(result_text or "", tmp_out)
                            with open(tmp_out, "rb") as fp:
                                buf = io.BytesIO(fp.read())
                            try:
                                if tmp_out and os.path.exists(tmp_out):
                                    os.remove(tmp_out)
                            except Exception:
                                pass
                        elif ext == ".xlsx":
                            from openpyxl import Workbook

                            wb = Workbook()
                            ws = wb.active
                            ws.title = "Tarjima"
                            for i, line in enumerate((result_text or "").splitlines(), start=1):
                                ws.cell(row=i, column=1, value=line)
                                if i >= 50000:
                                    break
                            bio = io.BytesIO()
                            wb.save(bio)
                            bio.seek(0)
                            buf = bio
                    except Exception:
                        buf = None

                elif act == "spellcheck":
                    try:
                        if ext == ".docx":
                            from docx import Document

                            doc = Document()
                            for line in (result_text or "").splitlines():
                                doc.add_paragraph(line)
                            bio = io.BytesIO()
                            doc.save(bio)
                            bio.seek(0)
                            buf = bio
                        elif ext == ".pptx":
                            from pptx import Presentation
                            from pptx.util import Inches

                            prs = Presentation()
                            slide = prs.slides.add_slide(prs.slide_layouts[5])
                            tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(5.0))
                            tf = tx.text_frame
                            for i, line in enumerate((result_text or "").splitlines()[:200]):
                                if i == 0:
                                    tf.text = line
                                else:
                                    tf.add_paragraph().text = line
                            bio = io.BytesIO()
                            prs.save(bio)
                            bio.seek(0)
                            buf = bio
                        elif ext == ".pdf":
                            fd_out, tmp_out = tempfile.mkstemp(suffix=".pdf", prefix="web_spell_notify_")
                            os.close(fd_out)
                            _build_pdf_from_text(result_text or "", tmp_out)
                            with open(tmp_out, "rb") as fp:
                                buf = io.BytesIO(fp.read())
                            try:
                                if tmp_out and os.path.exists(tmp_out):
                                    os.remove(tmp_out)
                            except Exception:
                                pass
                        elif ext == ".xlsx":
                            from openpyxl import Workbook

                            wb = Workbook()
                            ws = wb.active
                            ws.title = "Tuzatilgan"
                            for i, line in enumerate((result_text or "").splitlines(), start=1):
                                ws.cell(row=i, column=1, value=line)
                                if i >= 50000:
                                    break
                            bio = io.BytesIO()
                            wb.save(bio)
                            bio.seek(0)
                            buf = bio
                    except Exception:
                        buf = None

                elif act == "translit_auto":
                    try:
                        if ext == ".docx":
                            from docx import Document
                            from backend.services.auto_script import auto_cyrillic_latin as _auto
                            from bot.services.transliterate_service import transliterate as _trl

                            det = _auto(text or "")
                            if det.direction == "none":
                                raise Exception("direction none")
                            doc = Document(io.BytesIO(raw))
                            for para in doc.paragraphs:
                                for run in para.runs:
                                    if run.text:
                                        run.text = _trl(run.text, det.direction)
                            for table in doc.tables:
                                for row in table.rows:
                                    for cell in row.cells:
                                        for para in cell.paragraphs:
                                            for run in para.runs:
                                                if run.text:
                                                    run.text = _trl(run.text, det.direction)
                            bio = io.BytesIO()
                            doc.save(bio)
                            bio.seek(0)
                            buf = bio
                        elif ext == ".pptx":
                            from pptx import Presentation
                            from backend.services.auto_script import auto_cyrillic_latin as _auto
                            from bot.services.transliterate_service import transliterate as _trl

                            det = _auto(text or "")
                            if det.direction == "none":
                                raise Exception("direction none")
                            prs = Presentation(io.BytesIO(raw))
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if shape.has_text_frame:
                                        for para in shape.text_frame.paragraphs:
                                            for run in para.runs:
                                                if run.text:
                                                    run.text = _trl(run.text, det.direction)
                                    if shape.has_table:
                                        for row in shape.table.rows:
                                            for cell in row.cells:
                                                for para in cell.text_frame.paragraphs:
                                                    for run in para.runs:
                                                        if run.text:
                                                            run.text = _trl(run.text, det.direction)
                            bio = io.BytesIO()
                            prs.save(bio)
                            bio.seek(0)
                            buf = bio
                        elif ext == ".pdf":
                            fd_out, tmp_out = tempfile.mkstemp(suffix=".pdf", prefix="web_trl2_notify_")
                            os.close(fd_out)
                            _build_pdf_from_text(result_text or "", tmp_out)
                            with open(tmp_out, "rb") as fp:
                                buf = io.BytesIO(fp.read())
                            try:
                                if tmp_out and os.path.exists(tmp_out):
                                    os.remove(tmp_out)
                            except Exception:
                                pass
                        elif ext == ".xlsx":
                            from openpyxl import Workbook

                            wb = Workbook()
                            ws = wb.active
                            ws.title = "O'girildi"
                            for i, line in enumerate((result_text or "").splitlines(), start=1):
                                ws.cell(row=i, column=1, value=line)
                                if i >= 50000:
                                    break
                            bio = io.BytesIO()
                            wb.save(bio)
                            bio.seek(0)
                            buf = bio
                    except Exception:
                        buf = None

                if buf is None:
                    out_name = f"{base_name}_{act}.txt"
                    buf = io.BytesIO((result_text or "").encode("utf-8"))
                    buf.seek(0)
                else:
                    buf.seek(0)
                await ptb.bot.send_document(
                    chat_id=chat_id,
                    document=InputFile(buf, filename=out_name),
                    caption=cap,
                )
            except Exception:
                return

        try:
            asyncio.create_task(_notify_bg())
        except Exception:
            pass

    return {
        "ok": True,
        "action": act,
        "text": result_text,
        "meta": meta,
        "quota": quota,
        "storage": (
            {
                "bucket": st.bucket,
                "path": st.path,
                "public_url": st.public_url,
                "signed_url": signed,
            }
            if st
            else None
        ),
    }


@router.post("/api/files/process_async")
async def api_process_file_async(
    file: UploadFile = File(...),
    action: str = Form(..., description="spellcheck|translit_auto|translate_auto"),
    target_lang: Optional[str] = Form(None, description="uz|ru|en (for translate_auto)"),
    token: Optional[str] = Form(None),
    telegram_id: Optional[str] = Form(None),
):
    """
    Async-first file pipeline:
      upload -> storage -> enqueue Celery task -> return job_id
    Client polls /api/jobs/{job_id}.
    """
    uid = resolve_telegram_uid(telegram_id, token)
    if not uid:
        raise HTTPException(status_code=401, detail="Foydalanuvchi aniqlanmadi")
    uid_int = int(uid)

    act = (action or "").strip().lower()
    if act not in {"spellcheck", "translit_auto", "translate_auto"}:
        raise HTTPException(status_code=400, detail="action noto'g'ri")

    # Quota gate (fast fail before upload processing)
    cat = None
    consumed = False
    if act == "spellcheck":
        from bot.services.plan_limits import CAT_SPELL

        cat = CAT_SPELL
    elif act == "translit_auto":
        from bot.services.plan_limits import CAT_TRANSLIT

        cat = CAT_TRANSLIT
    else:
        from bot.services.plan_limits import CAT_TRANSLATE

        cat = CAT_TRANSLATE
    consumed = bool(web_quota_consume_or_raise(uid_int, cat))

    try:
        raw = await read_upload_limited(file)
    except EmptyUploadError:
        raise HTTPException(status_code=400, detail="Bo'sh fayl")
    except UploadTooLargeError as e:
        raise HTTPException(status_code=413, detail=str(e))

    fname = file.filename or "upload.txt"
    if "." not in fname:
        ct = (file.content_type or "").lower()
        if "wordprocessingml" in ct:
            fname = fname + ".docx"
        elif "pdf" in ct:
            fname = fname + ".pdf"
        elif "text" in ct:
            fname = fname + ".txt"

    st = upload_bytes_for_user(telegram_id=uid_int, filename=fname, raw=raw)
    if not st:
        # Without shared storage, we cannot safely async-process across services.
        raise HTTPException(status_code=503, detail="File storage is not configured. Try again later.")

    signed = create_signed_url(bucket=st.bucket, path=st.path, expires_in=3600)

    from backend.jobs import create_job
    from backend.celery_app import celery_app

    try:
        job = await create_job(kind=f"files:{act}")
    except Exception as e:
        # Most common cause: Redis not configured/reachable
        logger.error("async file queue unavailable: %s", e, exc_info=True)
        raise HTTPException(status_code=503, detail="Queue temporarily unavailable. Please try again later.")
    payload = {
        "job_id": job.job_id,
        "user_id": uid_int,
        "action": act,
        "target_lang": (target_lang or "").strip().lower(),
        "filename": fname,
        "storage": {"bucket": st.bucket, "path": st.path, "public_url": st.public_url, "signed_url": signed},
        "consumed": bool(consumed),
        "category": cat,
    }
    celery_app.send_task("files.process_v1", args=[payload], kwargs={}, task_id=job.job_id)

    from bot.services.plan_limits import category_quota_for_response
    from bot.services.user_service import record_service_completion

    # Record completion here as "enqueued" audit only; quota already consumed.
    # Final success/failure is tracked by job result.
    try:
        record_service_completion(uid_int, cat, f"Web async file:{act}", skip_quota=bool(consumed))
    except Exception:
        pass
    quota = None
    try:
        quota = category_quota_for_response(uid_int, cat)
    except Exception:
        quota = None

    return {
        "ok": True,
        "job_id": job.job_id,
        "status": job.status,
        "action": act,
        "quota": quota,
        "storage": {"bucket": st.bucket, "path": st.path, "public_url": st.public_url, "signed_url": signed},
    }

@router.post("/api/objective")
async def api_objective(req: ObjectiveRequest):
    role = (req.role or "").strip()
    if not role:
        raise HTTPException(status_code=400, detail="Role bo'sh bo'lishi mumkin emas")
    if len(role) > 120:
        raise HTTPException(status_code=400, detail="Role juda uzun")
    if req.extra and len(req.extra) > 800:
        raise HTTPException(status_code=400, detail="Qo'shimcha ma'lumot juda uzun")

    try:
        from bot.services.ai_service import generate_objective

        text = await generate_objective(
            role=role,
            experience=req.experience,
            extra=req.extra or "",
            lang=req.lang or "uz",
        )
        if not text or text.startswith("AI model"):
            raise HTTPException(status_code=502, detail=text or "AI javobi bo'sh")
        if text.startswith("Xatolik"):
            raise HTTPException(status_code=502, detail=text)
        try:
            uid = _resolve_web_uid_optional(req.telegram_id, req.token)
            if uid:
                from bot.utils.action_logger import log_action_fire_and_forget

                log_action_fire_and_forget(
                    telegram_id=int(uid),
                    username=None,
                    action_type="OBJECTIVE",
                    details=f"web:{(req.lang or 'uz')}",
                    metadata={"role": role[:120], "has_extra": bool(req.extra)},
                )
        except Exception:
            pass
        return {"ok": True, "text": text}
    except HTTPException:
        raise
    except Exception as e:
        logger.error("Objective API error: %s", e, exc_info=True)
        raise HTTPException(status_code=500, detail=f"Objective server xatosi: {str(e)[:200]}")


@router.post("/api/cv_parse")
async def api_cv_parse(
    file: UploadFile = File(...),
    token: Optional[str] = Form(None),
    telegram_id: Optional[str] = Form(None),
):
    """
    CV/Objectives parser:
      - Extract text from DOCX/PDF/PPTX/TXT/XLSX
      - Return structured result: name, skills, experience
    """
    uid = resolve_telegram_uid(telegram_id, token)
    uid_int = int(uid) if uid else None
    try:
        raw = await read_upload_limited(file)
    except EmptyUploadError:
        raise HTTPException(status_code=400, detail="Bo'sh fayl")
    except UploadTooLargeError as e:
        raise HTTPException(status_code=413, detail=str(e))

    fname = file.filename or "cv.txt"
    if "." not in fname:
        ct = (file.content_type or "").lower()
        if "wordprocessingml" in ct:
            fname += ".docx"
        elif "presentationml" in ct or "powerpoint" in ct:
            fname += ".pptx"
        elif "pdf" in ct:
            fname += ".pdf"
        elif "spreadsheetml" in ct or "excel" in ct:
            fname += ".xlsx"
        else:
            fname += ".txt"

    storage = None
    try:
        if uid_int:
            storage = upload_bytes_for_user(telegram_id=uid_int, filename=fname, raw=raw)
    except Exception:
        storage = None

    try:
        from bot.services.document_text_extract import extract_plain_text_from_bytes

        text = extract_plain_text_from_bytes(fname, raw)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e

    if not (text or "").strip():
        raise HTTPException(status_code=422, detail="Fayldan matn ajratilmadi")

    parsed = parse_cv_text(text)

    try:
        if uid_int:
            from bot.utils.action_logger import log_action_fire_and_forget

            log_action_fire_and_forget(
                telegram_id=int(uid_int),
                username=None,
                action_type="CV_PARSE",
                details="web:file",
                metadata={
                    "filename": fname,
                    "bytes": len(raw),
                    "chars": len(text or ""),
                    "skills": len(parsed.skills or []),
                    "experience": len(parsed.experience or []),
                    "has_name": bool(parsed.name),
                },
            )
    except Exception:
        pass

    return {
        "ok": True,
        "filename": fname,
        "name": parsed.name,
        "skills": parsed.skills,
        "experience": parsed.experience,
        "sections": parsed.raw_sections,
        "storage": (
            {
                "bucket": storage.bucket,
                "path": storage.path,
                "public_url": storage.public_url,
                "signed_url": create_signed_url(bucket=storage.bucket, path=storage.path, expires_in=3600),
            }
            if storage
            else None
        ),
    }
