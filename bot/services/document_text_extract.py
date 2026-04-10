"""Extract plain text from uploaded document bytes (WebApp spellcheck, APIs)."""
from __future__ import annotations

import io
import logging
from typing import Final

logger = logging.getLogger(__name__)

_SUPPORTED: Final = frozenset({"txt", "docx", "pptx", "pdf", "xlsx"})


def extract_plain_text_from_bytes(filename: str, raw: bytes) -> str:
    if not raw:
        raise ValueError("Empty file")
    ext = (filename or "").rsplit(".", 1)[-1].lower()
    if ext not in _SUPPORTED:
        raise ValueError(f"Unsupported type .{ext}; use .txt, .docx, .pptx, .pdf, .xlsx")

    if ext == "txt":
        for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
            try:
                return raw.decode(enc)
            except Exception:
                continue
        return raw.decode("utf-8", errors="replace")

    if ext == "docx":
        from docx import Document

        doc = Document(io.BytesIO(raw))
        parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append("\t".join(cells))
        for section in getattr(doc, "sections", []):
            for part in (getattr(section, "header", None), getattr(section, "footer", None)):
                if not part:
                    continue
                for p in getattr(part, "paragraphs", []):
                    t = (getattr(p, "text", "") or "").strip()
                    if t:
                        parts.append(t)
                for table in getattr(part, "tables", []):
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append("\t".join(cells))
        return "\n".join(parts)

    if ext == "pptx":
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        parts: list[str] = []
        def _walk_shapes(shapes):
            for shape in shapes:
                yield shape
                nested = getattr(shape, "shapes", None)
                if nested:
                    try:
                        yield from _walk_shapes(nested)
                    except Exception:
                        pass
        for slide in prs.slides:
            for shape in _walk_shapes(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = (cell.text or "").strip()
                            if t:
                                parts.append(t)
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide:
                for shape in _walk_shapes(notes_slide.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
        return "\n".join(parts)

    if ext == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError as e:
            logger.error("pypdf missing for PDF text extract: %s", e)
            raise ValueError("PDF imlo uchun serverda pypdf kutubxonasi kerak") from e
        reader = PdfReader(io.BytesIO(raw))
        texts: list[str] = []
        for page in reader.pages:
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                texts.append(t)
        return "\n".join(texts).strip()

    if ext == "xlsx":
        """
        Excel: best-effort text extraction (sheet by sheet).
        Keeps row/cell separation with tabs/newlines.
        """
        try:
            from openpyxl import load_workbook
        except Exception as e:
            logger.error("openpyxl missing for XLSX extract: %s", e)
            raise ValueError("XLSX uchun serverda openpyxl kutubxonasi kerak") from e
        try:
            wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        except Exception as e:
            raise ValueError("XLSX faylni o‘qib bo‘lmadi") from e
        parts: list[str] = []
        try:
            for ws in wb.worksheets:
                parts.append(f"--- {ws.title} ---")
                # Iterate rows; cap extremely wide rows
                for row in ws.iter_rows(values_only=True):
                    if not row:
                        continue
                    cells: list[str] = []
                    for v in row[:200]:
                        if v is None:
                            cells.append("")
                        else:
                            s = str(v).strip()
                            cells.append(s)
                    line = "\t".join(cells).strip()
                    if line:
                        parts.append(line)
        finally:
            try:
                wb.close()
            except Exception:
                pass
        return "\n".join(parts).strip()

    raise ValueError("Unsupported format")
