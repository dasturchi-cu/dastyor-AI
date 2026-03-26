"""
AI Service Module - Google Gemini Integration (Async)
Handles Text Processing, Data Extraction, and Translation using Gemini Asynchronously.
"""
import logging
import json
import asyncio
import os
import re
import httpx
import subprocess
import tempfile
import warnings
from difflib import SequenceMatcher
from difflib import get_close_matches
import shutil

# Eski paket: google-generativeai; keyinroq google-genai ga ko'chirish mumkin
warnings.filterwarnings(
    "ignore",
    message=".*google.generativeai.*",
    category=DeprecationWarning,
)

import google.generativeai as genai
from config import GOOGLE_API_KEY

logger = logging.getLogger(__name__)

# Initialize Gemini
if GOOGLE_API_KEY:
    try:
        genai.configure(api_key=GOOGLE_API_KEY)
        logger.info("Google Gemini initialized successfully")
    except Exception as e:
        logger.error(f"Failed to init Gemini: {e}")

# Models to try in order (newest first — keeps working as Gemini releases progress)
GEMINI_MODELS = [
    'gemini-2.5-flash',
    'gemini-2.0-flash',
    'gemini-1.5-flash-latest',
]

_MODEL_CACHE = None
_MODEL_CACHE_NAME = None
_MODEL_CACHE_LOCK = asyncio.Lock()

_LANGTOOL_RU = None
_LANGTOOL_RU_LOCK = asyncio.Lock()


async def get_model(preferred_models: list[str] | None = None):
    """Get Gemini model instance (cached) — tries models in order."""
    if not GOOGLE_API_KEY:
        return None

    global _MODEL_CACHE, _MODEL_CACHE_NAME
    order = preferred_models or GEMINI_MODELS

    # Fast path: cached model still matches preference list (or default list).
    if _MODEL_CACHE is not None and _MODEL_CACHE_NAME in order:
        return _MODEL_CACHE

    async with _MODEL_CACHE_LOCK:
        if _MODEL_CACHE is not None and _MODEL_CACHE_NAME in order:
            return _MODEL_CACHE

        for model_name in order:
            try:
                model = genai.GenerativeModel(model_name)
                _MODEL_CACHE = model
                _MODEL_CACHE_NAME = model_name
                logger.info(f"Using Gemini model: {model_name}")
                return model
            except Exception as e:
                logger.warning(f"Model {model_name} unavailable: {e}")

        logger.error("All Gemini models unavailable!")
        return None


GEMINI_TIMEOUT = 90  # seconds per API call
STT_FILE_PROCESS_MAX_WAIT_SECONDS = float(os.getenv("STT_FILE_PROCESS_MAX_WAIT_SECONDS", "45") or "45")
STT_FILE_PROCESS_POLL_SECONDS = float(os.getenv("STT_FILE_PROCESS_POLL_SECONDS", "0.5") or "0.5")
STT_SKIP_AUDIO_NORMALIZE = str(os.getenv("STT_SKIP_AUDIO_NORMALIZE", "0") or "0").strip().lower() in ("1", "true", "yes")

async def _gcall(coro, timeout: int = GEMINI_TIMEOUT):
    """Wrap generate_content_async with a hard timeout so the bot never hangs."""
    try:
        return await asyncio.wait_for(coro, timeout=timeout)
    except asyncio.TimeoutError:
        logger.error(f"Gemini API call timed out after {timeout}s")
        return None


def _set_para_text(para, text: str):
    """Set a Word paragraph's text safely (python-docx Paragraph has no .text setter)."""
    for run in para.runs:
        run.text = ''
    if para.runs:
        para.runs[0].text = text
    else:
        para.add_run(text)


def _set_pptx_paragraph_text(para, text: str):
    """
    Replace paragraph text while keeping at least one run style alive.
    """
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.add_run().text = text


_STT_FAILURE_SUBSTRINGS = (
    "Audio yuklashda xatolik",
    "qayta ishlashda xatolik",
    "AI model xatosi",
    "Audio tanilmadi",
    "Bo'sh javob",
    "transkripsiya xatoligi",
)


def is_valid_transcription_text(text: str) -> bool:
    """True agar matn haqiqiy transkripsiya bo'lsa (xato xabari emas)."""
    t = (text or "").strip()
    if not t:
        return False
    u = t.upper()
    if u in ("[EMPTY]", "[BO'SH]", "[BOSH]"):
        return False
    if any(s in t for s in _STT_FAILURE_SUBSTRINGS):
        return False
    # Kamida 2 ta "belgi" (harf yoki raqam) — qisqa "ha" / "yo'q" ham o'tsin
    sig = sum(1 for ch in t if ch.isalnum())
    return sig >= 2


def _normalize_audio_for_gemini(src_path: str) -> tuple[str, bool]:
    """
    Telegram ovozi ko'pincha OGG/Opus. Gemini yuklashini yaxshilash uchun MP3 ga aylantiramiz.
    Qaytaradi: (yuklash uchun yo'l, vaqtinchalik fayl bo'lsa True).
    """
    src_path = os.path.abspath(src_path)
    if not os.path.isfile(src_path):
        return src_path, False
    ext = (os.path.splitext(src_path)[1] or "").lower()
    if ext in (".mp3", ".wav", ".mpeg", ".mpga", ".flac"):
        return src_path, False
    if not shutil.which("ffmpeg"):
        logger.debug("ffmpeg yo'q — audio asl formatda yuklanadi")
        return src_path, False
    fd, out_path = tempfile.mkstemp(suffix=".mp3", prefix="stt_")
    os.close(fd)
    try:
        subprocess.run(
            [
                "ffmpeg",
                "-y",
                "-i",
                src_path,
                "-ac",
                "1",
                "-ar",
                "16000",
                "-codec:a",
                "libmp3lame",
                "-b:a",
                "64k",
                out_path,
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
        if os.path.isfile(out_path) and os.path.getsize(out_path) > 80:
            return out_path, True
    except Exception as e:
        logger.warning("ffmpeg audio normalize: %s", e)
    try:
        os.remove(out_path)
    except OSError:
        pass
    return src_path, False


async def transcribe_audio(audio_file_path: str) -> str:
    """
    Transcribe audio using Gemini 3 Flash (Multimodal) asynchronously
    """
    if not GOOGLE_API_KEY:
        return ""
    
    loop = asyncio.get_running_loop()
    upload_path, temp_conv = (
        (audio_file_path, False)
        if STT_SKIP_AUDIO_NORMALIZE
        else _normalize_audio_for_gemini(audio_file_path)
    )

    def blocking_upload(path: str):
        import time

        try:
            myfile = genai.upload_file(path)
            waited = 0.0
            while getattr(myfile.state, "name", str(myfile.state)) == "PROCESSING" and waited < STT_FILE_PROCESS_MAX_WAIT_SECONDS:
                time.sleep(STT_FILE_PROCESS_POLL_SECONDS)
                waited += STT_FILE_PROCESS_POLL_SECONDS
                try:
                    myfile = genai.get_file(myfile.name)
                except Exception:
                    break
            st = getattr(myfile.state, "name", str(myfile.state))
            if st == "PROCESSING":
                logger.warning("STT: fayl hali PROCESSING (timeout)")
            return myfile
        except Exception as e:
            logger.error(f"Upload error: {e}")
            return None

    try:
        logger.info("Uploading audio for STT: %s (normalized=%s)", upload_path, temp_conv)
        myfile = await loop.run_in_executor(None, blocking_upload, upload_path)
        
        if not myfile:
            return "Audio yuklashda xatolik."
            
        if myfile.state.name == "FAILED":
            return "Audio faylni qayta ishlashda xatolik."
        
        model = await get_model()
        if not model: return "AI model xatosi."
        
        result = await _gcall(model.generate_content_async(
            [
                myfile,
                (
                    "Transcribe the speech in this audio to plain text. "
                    "Preserve Uzbek (Latin or Cyrillic) as spoken; do not translate to English. "
                    "If the audio is silent or unintelligible, reply with exactly: [EMPTY]. "
                    "Otherwise output only the transcript, no labels or commentary."
                ),
            ]
        ))

        # Cleanup (non-blocking)
        async def cleanup():
            try:
                genai.delete_file(myfile.name)
            except Exception:
                logger.debug("genai.delete_file failed name=%s", getattr(myfile, "name", None), exc_info=True)
        asyncio.create_task(cleanup())

        if not result or not result.candidates:
            fb = getattr(result, "prompt_feedback", None) if result else None
            if fb:
                logger.warning("STT prompt_feedback: %s", fb)
            return "Audio tanilmadi."

        raw = (result.text or "").strip()
        if raw.upper() in ("[EMPTY]", "[BO'SH]", "[BOSH]"):
            return ""
        logger.info("STT raw len=%s preview=%r", len(raw), raw[:160])
        return raw if raw else "Bo'sh javob."

    except Exception as e:
        logger.error(f"Gemini Async Transcription error: {e}", exc_info=True)
        return "Audio transkripsiya xatoligi."
    finally:
        if temp_conv:
            try:
                os.remove(upload_path)
            except OSError:
                pass


async def translate_document_gemini(file_path: str, target_language: str = "uz") -> str:
    """
    Translates a document using Gemini AI asynchronously.
    """
    if not GOOGLE_API_KEY:
        return ""
        
    try:
        from docx import Document
        
        # Doc processing is CPU bound, run in Executor
        loop = asyncio.get_running_loop()
        doc = await loop.run_in_executor(None, Document, file_path)
        
        model = await get_model()
        
        lang_map = {"uz": "Uzbek", "ru": "Russian", "en": "English"}
        target_lang_name = lang_map.get(target_language, "Uzbek")
        
        # Collect every non-empty paragraph separately (better structure preservation).
        para_nodes = []
        for para in doc.paragraphs:
            if (para.text or "").strip():
                para_nodes.append(para)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if (para.text or "").strip():
                            para_nodes.append(para)

        logger.info(f"Translating doc: {len(para_nodes)} paragraphs to {target_lang_name}")

        async def translate_one(src_text: str) -> str:
            prompt = (
                f"Translate the following text to {target_lang_name}.\n"
                "Rules:\n"
                "- Return ONLY translated text.\n"
                "- Keep numbers, names, abbreviations, and punctuation.\n"
                "- Do not add explanations.\n\n"
                f"{src_text}"
            )
            try:
                resp = await _gcall(model.generate_content_async(prompt))
                out = (resp.text if resp and resp.text else "").strip()
                if out.startswith("```"):
                    out = out.replace("```text", "").replace("```", "").strip()
                return out or src_text
            except Exception as e:
                logger.error(f"Doc paragraph translation error: {e}")
                return src_text

        # Process paragraphs in moderate batches.
        for i in range(0, len(para_nodes), 8):
            batch_nodes = para_nodes[i:i + 8]
            srcs = [p.text for p in batch_nodes]
            outs = await asyncio.gather(*[translate_one(s) for s in srcs])
            for p, out in zip(batch_nodes, outs):
                _set_para_text(p, out)
            await asyncio.sleep(0.15)

        # Save
        output_path = file_path.replace(".docx", f"_translated_{target_language}.docx")
        await loop.run_in_executor(None, doc.save, output_path)
        return output_path

    except Exception as e:
        logger.error(f"Async Doc Translation failed: {e}", exc_info=True)
        return ""


async def translate_text(text: str, direction: str = "uz_en") -> str:
    """
    Translate plain text using Gemini. Used by /api/translate web endpoint.
    direction: uz_en | en_uz | ru_uz | uz_ru | ru_en | en_ru
    """
    lang_map = {
        "uz_en": ("O'zbek", "English"),
        "en_uz": ("English", "O'zbek"),
        "ru_uz": ("Russian", "O'zbek"),
        "uz_ru": ("O'zbek", "Russian"),
        "ru_en": ("Russian", "English"),
        "en_ru": ("English", "Russian"),
    }
    src, tgt = lang_map.get(direction, ("O'zbek", "English"))
    code_map = {
        "uz_en": ("uz", "en"),
        "en_uz": ("en", "uz"),
        "ru_uz": ("ru", "uz"),
        "uz_ru": ("uz", "ru"),
        "ru_en": ("ru", "en"),
        "en_ru": ("en", "ru"),
    }
    sl, tl = code_map.get(direction, ("uz", "en"))

    async def _google_fallback_translate(src_text: str, source_lang: str, target_lang: str) -> str:
        # Public translate endpoint fallback to keep service alive when Gemini quota is exhausted.
        url = "https://translate.googleapis.com/translate_a/single"
        params = {
            "client": "gtx",
            "sl": source_lang,
            "tl": target_lang,
            "dt": "t",
            "q": src_text,
        }
        timeout = httpx.Timeout(12.0)
        async with httpx.AsyncClient(timeout=timeout) as client:
            resp = await client.get(url, params=params)
            resp.raise_for_status()
            data = resp.json()
        parts = []
        if isinstance(data, list) and data and isinstance(data[0], list):
            for item in data[0]:
                if isinstance(item, list) and item:
                    parts.append(str(item[0] or ""))
        return "".join(parts).strip()

    model = await get_model()

    prompt = (
        f"Translate the following {src} text to {tgt}.\n"
        "Return ONLY the translated text, no explanations.\n\n"
        f"{text}"
    )
    try:
        if model:
            resp = await _gcall(model.generate_content_async(prompt))
            out = resp.text.strip() if resp and resp.text else ""
        else:
            out = ""
        if not out:
            out = await _google_fallback_translate(text, sl, tl)
        if not is_meaningfully_changed(text, out):
            return "Tarjima natijasi original bilan bir xil chiqdi."
        return out
    except Exception as e:
        logger.warning("Gemini translate failed, using fallback: %s", e)
        try:
            out = await _google_fallback_translate(text, sl, tl)
            if not out:
                return "Tarjima vaqtincha mavjud emas."
            if not is_meaningfully_changed(text, out):
                return "Tarjima natijasi original bilan bir xil chiqdi."
            return out
        except Exception as e2:
            logger.error("translate_text fallback error: %s", e2, exc_info=True)
            return "Tarjima vaqtincha mavjud emas."


async def translate_pptx(file_path: str, direction: str = "uz_en", target_language: str = "uz") -> str:
    """
    Translate PPTX text in-place and save as translated PPTX.
    Keeps slide layout, only replacing paragraph text.
    """
    try:
        from pptx import Presentation

        loop = asyncio.get_running_loop()
        prs = await loop.run_in_executor(None, Presentation, file_path)

        def _walk_shapes(shapes):
            for shape in shapes:
                yield shape
                nested = getattr(shape, "shapes", None)
                if nested:
                    try:
                        yield from _walk_shapes(nested)
                    except Exception:
                        pass

        paragraphs = []
        for slide in prs.slides:
            for shape in _walk_shapes(slide.shapes):
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if (para.text or "").strip():
                            paragraphs.append(para)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                if (para.text or "").strip():
                                    paragraphs.append(para)

        for i in range(0, len(paragraphs), 6):
            batch = paragraphs[i:i + 6]
            srcs = [p.text for p in batch]
            outs = await asyncio.gather(*[translate_text(s, direction) for s in srcs])
            for para, out in zip(batch, outs):
                if out and not out.startswith("Tarjima vaqtincha mavjud emas."):
                    _set_pptx_paragraph_text(para, out)
            await asyncio.sleep(0.1)

        output_path = file_path.replace(".pptx", f"_translated_{target_language}.pptx")
        await loop.run_in_executor(None, prs.save, output_path)
        return output_path
    except Exception as e:
        logger.error("translate_pptx failed: %s", e, exc_info=True)
        return ""


def count_words_text(text: str) -> int:
    """Whitespace-delimited tokens (works for Uzbek apostrophes and Cyrillic)."""
    s = (text or "").strip()
    if not s:
        return 0
    return len(re.findall(r"\S+", s))


def _normalize_for_compare(text: str) -> str:
    s = (text or "").strip().lower()
    # Normalize common apostrophe variants and whitespace.
    s = (
        s.replace("’", "'")
        .replace("`", "'")
        .replace("ʻ", "'")
        .replace("ʼ", "'")
        .replace("ʹ", "'")
    )
    s = re.sub(r"\s+", " ", s)
    return s


def is_meaningfully_changed(source: str, candidate: str) -> bool:
    a = _normalize_for_compare(source)
    b = _normalize_for_compare(candidate)
    if not a and not b:
        return False
    return a != b


async def check_spelling_text(text: str) -> tuple[str, int]:
    """
    Spell-check plain text (Uzbek/Russian) with conservative corrections.
    Returns: (corrected_text, fixes_count)
    """
    src = (text or "").strip()
    if not src:
        return "", 0
    original_src = src
    heuristic_fixes = 0

    # If there is nothing to fix (no letters), skip spell checking.
    if not any(ch.isalpha() for ch in src):
        return src, 0

    def _is_cyrillic(s: str) -> bool:
        # Rough heuristic: detect Cyrillic blocks.
        return any(("\u0400" <= ch <= "\u04FF") or ("\u0500" <= ch <= "\u052F") for ch in (s or ""))

    def _apply_uz_latin_heuristics(s: str) -> tuple[str, int]:
        """
        Fast deterministic fixes for common Uzbek-latin typos.
        """
        out = s
        rules = [
            (r"\bbugn\b", "bugun"),
            (r"\bozimni\b", "o'zimni"),
            (r"\bxis\b", "his"),
            (r"\bqildm\b", "qildim"),
            (r"\boqip\b", "o'qib"),
            (r"\bbolgach\b", "bo'lgach"),
            (r"\bqoydi\b", "qo'ydi"),
            (r"\bdstlarim\b", "do'stlarim"),
            (r"\bkochada\b", "ko'chada"),
            (r"\boynab\b", "o'ynab"),
            (r"\byurganmda\b", "yurganimda"),
            (r"\byomgr\b", "yomg'ir"),
            (r"\byogib\b", "yog'ib"),
        ]
        changes = 0
        for pat, rep in rules:
            new_out, n = re.subn(pat, rep, out, flags=re.IGNORECASE)
            if n:
                out = new_out
                changes += int(n)
        return out, changes

    uz_typo_map = {
        "bugn": "bugun",
        "bormadimn": "bormadim",
        "maktabgaa": "maktabga",
        "qacchon": "qachon",
        "xam": "ham",
        "bilanam": "bilan ham",
    }

    uz_dictionary = {
        "men", "sen", "u", "biz", "siz", "ular", "bugun", "kecha", "ertaga",
        "maktabga", "maktab", "bormadim", "bordim", "boraman", "kelaman", "keldim",
        "uyga", "ishga", "kitob", "daftar", "yozdim", "o'qidim", "dars", "darsga",
        "rahmat", "iltimos", "qanday", "yaxshi", "yomon", "bo'ldi", "bo'lmadi",
        "do'stim", "do'stlarim", "o'zim", "o'zimni", "o'qib", "ko'chada", "o'ynab",
        "yurganimda", "yomg'ir", "yog'ib", "uchun", "haqiqiy", "natija", "xato",
        "tekshiruv", "imlo", "tarjima", "matn", "fayl",
    }

    def _preserve_case(src_word: str, rep_word: str) -> str:
        if src_word.isupper():
            return rep_word.upper()
        if src_word[:1].isupper():
            return rep_word[:1].upper() + rep_word[1:]
        return rep_word

    def _apply_uz_dictionary_pass(s: str) -> tuple[str, int]:
        """
        Lightweight Uzbek latin spell pass with typo map + close match.
        Conservative by design to avoid semantic rewrites.
        """
        pattern = re.compile(r"[A-Za-zO'oʻʼ`’ʻ]+", flags=re.UNICODE)
        changes = 0

        def repl(m: re.Match) -> str:
            nonlocal changes
            word = m.group(0)
            wnorm = _normalize_for_compare(word)
            # Fast exact-known typo fixes.
            direct = uz_typo_map.get(wnorm)
            if direct:
                changes += 1
                return _preserve_case(word, direct)
            # Skip short tokens and known dictionary words.
            if len(wnorm) < 4 or wnorm in uz_dictionary:
                return word
            # Conservative nearest dictionary lookup.
            cands = get_close_matches(wnorm, uz_dictionary, n=1, cutoff=0.86)
            if not cands:
                return word
            best = cands[0]
            # Avoid aggressive replacements for much longer/shorter words.
            if abs(len(best) - len(wnorm)) > 2:
                return word
            changes += 1
            return _preserve_case(word, best)

        out = pattern.sub(repl, s)
        return out, changes

    def _mask_sensitive_tokens(s: str) -> tuple[str, dict[str, str]]:
        """
        Replace URLs/emails/numbers with placeholders to prevent the model
        from "correcting" them into something else.
        """
        token_re = re.compile(r"(https?://\S+|www\.\S+|\b\S+@\S+\b|\b\d+(?:[.,]\d+)?\b)")
        mapping: dict[str, str] = {}
        i = 0

        def repl(m: re.Match) -> str:
            nonlocal i
            key = f"__TOK{i}__"
            mapping[key] = m.group(0)
            i += 1
            return key

        masked = token_re.sub(repl, s)
        return masked, mapping

    def _unmask_sensitive_tokens(s: str, mapping: dict[str, str]) -> str:
        out = s
        for key, value in mapping.items():
            out = out.replace(key, value)
        return out

    def _conservative_merge(source: str, candidate: str) -> str:
        """
        Apply only small, local changes from `candidate` into `source`.
        We ignore pure insertions (to reduce hallucinations) and require
        similarity/length closeness for "replace" blocks.
        """
        s = source or ""
        c = (candidate or "")
        if not c.strip():
            return s

        # Quick divergence guard.
        if SequenceMatcher(None, s, c).ratio() < 0.6:
            return s

        sm = SequenceMatcher(None, s, c, autojunk=False)
        out: list[str] = []

        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            if tag == "equal":
                out.append(s[i1:i2])
                continue

            if tag == "replace":
                s_seg = s[i1:i2]
                c_seg = c[j1:j2]
                if not c_seg:
                    out.append(s_seg)
                    continue

                seg_ratio = SequenceMatcher(None, s_seg, c_seg).ratio()
                max_len = max(1, len(s_seg), len(c_seg))
                len_close = abs(len(s_seg) - len(c_seg)) / max_len <= 0.35

                if seg_ratio >= 0.78 and len_close:
                    out.append(c_seg)
                else:
                    out.append(s_seg)
                continue

            if tag == "delete":
                # Keep original text if model removed something.
                out.append(s[i1:i2])
                continue

            if tag == "insert":
                ins = c[j1:j2]
                if not ins:
                    continue
                # Allow insertion of whitespace/punctuation only.
                if ins.strip() == "":
                    out.append(ins)
                elif len(ins) <= 5 and not any(ch.isalnum() for ch in ins):
                    out.append(ins)
                elif len(ins) <= 3 and ins.isalpha():
                    # Spellcheckda eng ko'p holat: so'z ichiga 1-3 ta harf yetishmaydi
                    # (masalan: "beramdi" -> "beradimi"). Buni ruxsat beramiz,
                    # lekin faqat so'z ichida bo'lsa.
                    prev_ch = s[i1 - 1] if i1 > 0 else ""
                    next_ch = s[i1] if i1 < len(s) else ""
                    if prev_ch.isalpha() and next_ch.isalpha():
                        out.append(ins)
                else:
                    # Ignore letter/number insertions.
                    pass
                continue

        merged = "".join(out)
        # Final sanity guard.
        merged_ratio = SequenceMatcher(None, s, merged).ratio()
        if merged_ratio < 0.55:
            return s
        # Agar konservativ merge juda ko'p narsani kesib yuborsa, lekin candidate juda yaqin bo'lsa,
        # candidate'ni qabul qilamiz (real typo-fixlar yo'qolib ketmasin).
        cand_ratio = SequenceMatcher(None, s, c).ratio()
        if merged == s and cand_ratio >= 0.85 and abs(len(c) - len(s)) <= max(8, int(len(s) * 0.25)):
            return c
        return merged

    def _count_fixes(source: str, corrected: str) -> int:
        if corrected == source:
            return 0
        a = (source or "").split()
        b = (corrected or "").split()
        diff = sum(1 for i in range(min(len(a), len(b))) if a[i] != b[i]) + abs(len(a) - len(b))
        return max(1, min(diff, 999)) if diff > 0 else 1

    try:
        # 0) Quick deterministic Uzbek-latin typo pass.
        if not _is_cyrillic(src):
            htxt, hfix = _apply_uz_latin_heuristics(src)
            if hfix > 0:
                src = htxt
                heuristic_fixes += int(hfix)
            dtext, dfix = _apply_uz_dictionary_pass(src)
            if dfix > 0:
                src = dtext
                heuristic_fixes += int(dfix)

        # 1) Deterministic spell correction for Russian (if LanguageTool is available).
        # This reduces wrong “LLM hallucination” corrections for ru text.
        if _is_cyrillic(src):
            try:
                global _LANGTOOL_RU
                import language_tool_python  # optional dependency
                async with _LANGTOOL_RU_LOCK:
                    if _LANGTOOL_RU is None:
                        # Tool initialization (downloads LanguageTool) is blocking.
                        _LANGTOOL_RU = await asyncio.to_thread(language_tool_python.LanguageTool, "ru")

                    matches = await asyncio.to_thread(_LANGTOOL_RU.check, src)
                if not matches:
                    return src, 0
                # Correct needs the same lock for safety.
                async with _LANGTOOL_RU_LOCK:
                    corrected = await asyncio.to_thread(_LANGTOOL_RU.correct, src)
                merged = _conservative_merge(src, corrected)
                return merged, max(int(heuristic_fixes), _count_fixes(original_src, merged))
            except Exception as e:
                logger.warning(f"LanguageTool ru failed, falling back to Gemini: {e}")

        # 2) Gemini fallback for Uzbek and/or when LanguageTool fails.
        # Spellcheck should prioritize speed: try fastest models first.
        model = await get_model(preferred_models=[
            "gemini-2.5-flash",
            "gemini-2.0-flash",
            "gemini-1.5-flash-latest",
        ])
        if not model:
            # No model and no local spell checker: best-effort (no changes).
            return src, int(heuristic_fixes)

        spell_timeout = int((globals().get("GEMINI_SPELLCHECK_TIMEOUT") or 30))

        def _make_prompt(s: str) -> str:
            return (
                "You are a spelling corrector.\n"
                "Fix ONLY obvious spelling typos (including casing and punctuation spacing if needed).\n"
                "Do NOT rewrite sentences, change meaning, add new words, or remove existing ones.\n"
                "Keep line breaks as in the input.\n"
                "Return ONLY the corrected text.\n\n"
                f"{s}"
            )

        def _correct_part(part_src: str) -> str:
            masked, mapping = _mask_sensitive_tokens(part_src)
            # Call site runs in async context; this function just formats.
            return masked, mapping

        async def _call_gemini(masked_part: str) -> str:
            resp = await _gcall(
                model.generate_content_async(
                    _make_prompt(masked_part),
                    generation_config={"temperature": 0.0},
                ),
                timeout=spell_timeout,
            )
            txt = (resp.text if resp and resp.text else "") or ""
            # Safety: remove wrappers if model returns formatted block.
            txt = txt.replace("```text", "").replace("```", "").strip()
            if len(txt) >= 2 and txt[0] == txt[-1] and txt[0] in ("'", '"'):
                txt = txt[1:-1].strip()
            return txt

        # For long texts, split into chunks and process concurrently (larger batches = fewer round-trips).
        if len(src) > 2000:
            parts: list[str] = []
            buf: list[str] = []
            cur = 0
            for para in src.splitlines():
                p = para.rstrip()
                # preserve blank lines
                if not p:
                    if buf:
                        parts.append("\n".join(buf).strip())
                        buf, cur = [], 0
                    parts.append("")
                    continue
                if cur + len(p) + 1 > 2000 and buf:
                    parts.append("\n".join(buf).strip())
                    buf, cur = [p], len(p)
                else:
                    buf.append(p)
                    cur += len(p) + 1
            if buf:
                parts.append("\n".join(buf).strip())

            out_parts: list[str] = []
            batch: list[tuple[str, str, dict[str, str]]] = []

            # (orig_text, masked_text, mapping)
            for part in parts:
                if part == "":
                    out_parts.append("")
                    continue
                masked, mapping = _mask_sensitive_tokens(part)
                batch.append((part, masked, mapping))
                if len(batch) >= 8:
                    resps = await asyncio.gather(*[
                        _call_gemini(masked_part) for (_, masked_part, _) in batch
                    ])
                    for (orig, _, mapping), resp_text in zip(batch, resps):
                        candidate = _unmask_sensitive_tokens(resp_text.strip(), mapping) if resp_text else orig
                        out_parts.append(_conservative_merge(orig, candidate))
                    batch = []

            if batch:
                resps = await asyncio.gather(*[
                    _call_gemini(masked_part) for (_, masked_part, _) in batch
                ])
                for (orig, _, mapping), resp_text in zip(batch, resps):
                    candidate = _unmask_sensitive_tokens(resp_text.strip(), mapping) if resp_text else orig
                    out_parts.append(_conservative_merge(orig, candidate))

            corrected = "\n".join(out_parts).strip()
        else:
            masked, mapping = _mask_sensitive_tokens(src)
            resp_text = await _call_gemini(masked)
            candidate = _unmask_sensitive_tokens(resp_text.strip(), mapping) if resp_text else src
            corrected = _conservative_merge(src, candidate)

        fixes = max(int(heuristic_fixes), _count_fixes(original_src, corrected))
        return corrected, fixes

    except Exception as e:
        logger.error(f"check_spelling_text error: {e}", exc_info=True)
        fallback_fixes = max(int(heuristic_fixes), _count_fixes(original_src, src))
        return src, int(fallback_fixes)


async def generate_objective(role: str, experience: str = "junior", extra: str = "", lang: str = "uz") -> str:
    """
    Generate a short professional CV objective/summary.
    lang: uz | ru | en
    experience: junior | middle | senior | lead
    """
    model = await get_model()
    if not model:
        return "AI model mavjud emas."

    l = (lang or "uz").lower().strip()
    if l not in ("uz", "ru", "en"):
        l = "uz"

    exp_map = {
        "junior": {"uz": "boshlang'ich", "ru": "начальный", "en": "junior"},
        "middle": {"uz": "o'rta", "ru": "средний", "en": "mid-level"},
        "senior": {"uz": "katta mutaxassis", "ru": "старший специалист", "en": "senior"},
        "lead":   {"uz": "rahbar", "ru": "руководитель", "en": "lead"},
    }
    exp_key = (experience or "junior").lower().strip()
    exp_label = exp_map.get(exp_key, exp_map["junior"])[l]

    # Keep it short and usable: 2–3 sentences, no bullet lists.
    if l == "ru":
        prompt = (
            "Сгенерируй краткую профессиональную цель (objective) для резюме.\n"
            "Требования: 2–3 предложения, без списков, без лишних объяснений, деловой стиль.\n"
            f"Должность: {role}\n"
            f"Уровень: {exp_label}\n"
            f"Дополнительно (если есть): {extra}\n"
        )
    elif l == "en":
        prompt = (
            "Generate a short professional resume objective.\n"
            "Requirements: 2–3 sentences, no bullet points, no extra explanation, professional tone.\n"
            f"Role: {role}\n"
            f"Level: {exp_label}\n"
            f"Extra (if any): {extra}\n"
        )
    else:
        prompt = (
            "Rezyume uchun qisqa professional obyektivka (objective) yozing.\n"
            "Talablar: 2–3 gap, ro'yxatsiz, ortiqcha izohsiz, rasmiy uslub.\n"
            f"Kasb/Lavozim: {role}\n"
            f"Daraja: {exp_label}\n"
            f"Qo'shimcha (bo'lsa): {extra}\n"
        )

    try:
        resp = await _gcall(model.generate_content_async(prompt))
        if resp is None:
            return "AI javobi kechikdi. Iltimos, qayta urinib ko'ring."
        return resp.text.strip() if resp.text else "Natija bo'sh."
    except Exception as e:
        logger.error(f"generate_objective error: {e}", exc_info=True)
        return f"Xatolik: {str(e)[:200]}"


def _oby_extract_has_content(d: dict) -> bool:
    """Hech bo'lmaganda bitta mazmunli maydon bor-yo'qligini tekshirish."""
    if not d:
        return False
    skip = {"", "yo'q", "yoq", "yok", "нет", "no", "none"}
    for k, v in d.items():
        if k == "relatives" and isinstance(v, list):
            for r in v:
                if isinstance(r, dict) and any(
                    str(x).strip() for x in r.values() if x is not None
                ):
                    return True
            continue
        if k == "work_experience" and isinstance(v, list):
            for w in v:
                if not isinstance(w, dict):
                    continue
                for x in w.values():
                    s = str(x).strip().lower() if x is not None else ""
                    if s and s not in skip:
                        return True
            continue
        if isinstance(v, str):
            s = v.strip().lower()
            if s and s not in skip:
                return True
    return False


async def extract_obyektivka_data(text: str) -> dict:
    """
    Extract structured data from text using Gemini asynchronously
    """
    model = await get_model(preferred_models=[
        "gemini-2.0-flash",
        "gemini-1.5-flash-latest",
        "gemini-2.5-flash",
    ])
    if not model: return {}
    
    prompt = f"""
    Quyidagi matndan obyektivka uchun faktlarni JSONga ajrat.
    Qoidalar:
    - Faqat aniq aytilgan ma'lumotni joylashtir, qolgani bo'sh.
    - fullname: "Familiya Ism Sharif" (imkon qadar).
    - work_experience va relatives ro'yxat bo'lsin.
    - Javob faqat JSON bo'lsin (markdownsiz).

    Matn: {text}

    Faqat bitta JSON qaytar (markdown, ``` yo'q):
    {{
        "fullname": "Familiya Ism Sharif",
        "birthdate": "KK.OO.YYYY yoki bo'sh",
        "birthplace": "Viloyat, Tuman",
        "nation": "Millati",
        "party": "Partiyaviyligi",
        "education": "Ma'lumoti",
        "graduated": "Tamomlagan joyi va yili",
        "specialty": "Mutaxassisligi",
        "degree": "Ilmiy darajasi",
        "scientific_title": "Ilmiy unvoni",
        "languages": "Tillar",
        "military_rank": "Harbiy unvoni",
        "awards": "Mukofotlari",
        "deputy": "Deputatligi",
        "work_experience": [{{"year": "Yillar", "position": "Ish joyi"}}],
        "relatives": [{{"degree": "Qarindoshligi", "fullname": "F.I.SH", "birth_year_place": "Tug'ilgan yili va joyi", "work_place": "Ish joyi", "address": "Yashash manzili"}}]
    }}
    """
    
    try:
        response = await _gcall(model.generate_content_async(prompt), timeout=35)
        if not response or not response.text:
            return {}
        cleaned = response.text.replace('```json', '').replace('```', '').strip()
        start = cleaned.find('{')
        end = cleaned.rfind('}') + 1
        if start != -1 and end != -1:
            cleaned = cleaned[start:end]
        parsed = json.loads(cleaned)
        if isinstance(parsed, dict) and _oby_extract_has_content(parsed):
            return parsed
        return {}

    except Exception as e:
        logger.error(f"Async Data extraction error: {e}")
        return {}


async def check_spelling_gemini(file_path: str) -> tuple[str, int, int]:
    """
    Checks spelling in a DOCX file using Gemini asynchronously.
    Returns: (output_path, errors_found, errors_fixed)
    """
    try:
        from docx import Document

        loop = asyncio.get_running_loop()
        doc = await loop.run_in_executor(None, Document, file_path)

        output_path = file_path.replace(".docx", "_checked.docx")

        paragraphs_to_check = []
        seen: set[int] = set()

        def _add_para(para):
            if not para:
                return
            t = getattr(para, "text", "") or ""
            if not t.strip():
                return
            pid = id(para)
            if pid in seen:
                return
            seen.add(pid)
            paragraphs_to_check.append(para)

        # Main body paragraphs
        for para in doc.paragraphs:
            _add_para(para)

        # Body tables
        for table in getattr(doc, "tables", []):
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        _add_para(para)

        # Headers/footers (often contain lots of text)
        for section in getattr(doc, "sections", []):
            header = getattr(section, "header", None)
            footer = getattr(section, "footer", None)
            for part in (header, footer):
                if not part:
                    continue
                for para in getattr(part, "paragraphs", []):
                    _add_para(para)
                for table in getattr(part, "tables", []):
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                _add_para(para)

        errors_fixed = 0

        async def process_para(para):
            src = para.text or ""
            corrected, fixes = await check_spelling_text(src)
            if corrected and corrected != src:
                _set_para_text(para, corrected)
                return fixes if fixes > 0 else 1
            return 0

        batch_size = 8
        for i in range(0, len(paragraphs_to_check), batch_size):
            batch = paragraphs_to_check[i:i + batch_size]
            errors_fixed += sum(await asyncio.gather(*[process_para(p) for p in batch]))
            await asyncio.sleep(0.15)

        await loop.run_in_executor(None, doc.save, output_path)

        return output_path, errors_fixed, errors_fixed

    except Exception as e:
        logger.error(f"Async Spell check failed: {e}", exc_info=True)
        # Best-effort: still return a copied file so the bot can send it back.
        try:
            output_path = file_path.replace(".docx", "_checked.docx")
            shutil.copyfile(file_path, output_path)
            return output_path, 0, 0
        except Exception:
            return "", 0, 0


async def check_spelling_pptx(file_path: str) -> tuple[str, int, int]:
    """
    Checks spelling in a PPTX file using Gemini asynchronously.
    Iterates slides → shapes → text frames → paragraphs → runs.
    Returns: (output_path, errors_found, errors_fixed)
    """
    try:
        from pptx import Presentation

        loop = asyncio.get_running_loop()
        prs = await loop.run_in_executor(None, Presentation, file_path)

        output_path = file_path.replace(".pptx", "_checked.pptx")

        paragraphs_to_check = []
        seen: set[int] = set()

        def _add_para(para):
            if not para:
                return
            t = getattr(para, "text", "") or ""
            if not t.strip():
                return
            pid = id(para)
            if pid in seen:
                return
            seen.add(pid)
            paragraphs_to_check.append(para)

        def _walk_shapes(shapes):
            # Recursively traverse group shapes to avoid missing text.
            for shape in shapes:
                yield shape
                nested = getattr(shape, "shapes", None)
                if nested:
                    try:
                        yield from _walk_shapes(nested)
                    except Exception:
                        pass

        for slide in prs.slides:
            for shape in _walk_shapes(slide.shapes):
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        _add_para(para)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                _add_para(para)

            # Notes slide text (often ignored but can contain user content)
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide:
                for shape in _walk_shapes(notes_slide.shapes):
                    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            _add_para(para)
                    if getattr(shape, "has_table", False) and shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for para in cell.text_frame.paragraphs:
                                    _add_para(para)

        errors_fixed = 0

        async def process_para(para):
            src = para.text or ""
            corrected, fixes = await check_spelling_text(src)
            if corrected and corrected != src:
                _set_pptx_paragraph_text(para, corrected)
                return fixes if fixes > 0 else 1
            return 0

        batch_size = 8
        for i in range(0, len(paragraphs_to_check), batch_size):
            batch = paragraphs_to_check[i:i + batch_size]
            errors_fixed += sum(await asyncio.gather(*[process_para(p) for p in batch]))
            await asyncio.sleep(0.15)

        await loop.run_in_executor(None, prs.save, output_path)
        return output_path, errors_fixed, errors_fixed

    except Exception as e:
        logger.error(f"PPTX Spell check failed: {e}", exc_info=True)
        try:
            output_path = file_path.replace(".pptx", "_checked.pptx")
            shutil.copyfile(file_path, output_path)
            return output_path, 0, 0
        except Exception:
            return "", 0, 0
